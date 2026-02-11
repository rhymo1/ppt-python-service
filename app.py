from flask import Flask, request, jsonify, send_file
import zipfile
import xml.etree.ElementTree as ET
import os
import tempfile
import base64
from datetime import datetime
from io import BytesIO
from pptx import Presentation
from PIL import Image
import PyPDF2
import json
import re

app = Flask(__name__)

# Root route
@app.route('/')
def home():
    return jsonify({
        "status": "Python PPT Service is running",
        "version": "3.0",
        "endpoints": [
            "/health",
            "/extract",
            "/calculate",
            "/calculate-financial",
            "/generate",
            "/generate-charts"
        ]
    })

# Health check route
@app.route('/health')
def health():
    return jsonify({
        "status": "healthy",
        "timestamp": datetime.now().isoformat()
    })

# Helper: Extract text from PDF
def extract_text_from_pdf(pdf_bytes):
    """Extract all text from a PDF file"""
    try:
        pdf_file = BytesIO(pdf_bytes)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()
        
        return text
    except Exception as e:
        return f"Error extracting PDF: {str(e)}"

# Helper: Extract data from sqproj
def extract_from_sqproj(sqproj_bytes):
    """Extract data from .sqproj file (ZIP archive with XML/SQLite)"""
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.sqproj') as tmp:
            tmp.write(sqproj_bytes)
            tmp_path = tmp.name
        
        extracted = {}
        
        try:
            with zipfile.ZipFile(tmp_path, 'r') as zip_ref:
                for filename in zip_ref.namelist():
                    if filename.endswith('.xml'):
                        xml_content = zip_ref.read(filename)
                        try:
                            root = ET.fromstring(xml_content)
                            for elem in root.iter():
                                if elem.text and elem.text.strip():
                                    key = elem.tag.lower().replace(' ', '_')
                                    extracted[key] = elem.text.strip()
                        except:
                            pass
        except Exception as e:
            extracted["sqproj_error"] = str(e)
        
        os.unlink(tmp_path)
        
        return extracted
    
    except Exception as e:
        return {"error": f"Failed to process sqproj: {str(e)}"}

# Helper: Parse building data from PDF text
def parse_building_data(text):
    """Extract key building data from PDF text"""
    data = {}
    
    # Extract building address
    address_match = re.search(r'Gebäudeadresse[:\s]+([^\n]+)', text, re.IGNORECASE)
    if address_match:
        data['gebaeude_adresse'] = address_match.group(1).strip()
    
    # Extract year
    year_match = re.search(r'Baujahr[:\s]+(\d{4})', text, re.IGNORECASE)
    if year_match:
        data['baujahr'] = year_match.group(1)
    
    # Extract building type
    type_match = re.search(r'Gebäudetyp[:\s]+([^\n]+)', text, re.IGNORECASE)
    if type_match:
        data['gebaeude_typ'] = type_match.group(1).strip()
    
    # Extract living area
    area_match = re.search(r'Wohnfläche[:\s]+ca\.\s*([\d,\.]+)\s*m', text, re.IGNORECASE)
    if area_match:
        data['wohnflaeche'] = area_match.group(1).strip()
    
    # Extract energy costs
    costs_match = re.search(r'Energiekosten[:\s]+([\d,\.]+)\s*€', text, re.IGNORECASE)
    if costs_match:
        data['energiekosten_aktuell'] = costs_match.group(1).strip()
    
    # Extract CO2 emissions
    co2_match = re.search(r'CO[²2]-Emission[:\s]+([\d,]+)\s*kg', text, re.IGNORECASE)
    if co2_match:
        data['co2_emission_aktuell'] = co2_match.group(1).strip()
    
    return data

# Extract data from files
@app.route('/extract', methods=['POST'])
def extract_data():
    try:
        extracted_data = {}
        processing_log = []
        files_processed = {
            "sqproj": False,
            "sanierungsfahrplan_pdf": False,
            "umsetzungshilfe_pdf": False
        }
        
        # Check if we're receiving files via multipart/form-data
        if request.files and len(request.files) > 0:
            processing_log.append("Receiving files via multipart/form-data")
            
            # Process each uploaded file
            for key in request.files:
                file = request.files[key]
                filename = file.filename.lower()
                file_content = file.read()
                
                processing_log.append(f"Processing file: {file.filename} ({len(file_content)} bytes)")
                
                # Detect .sqproj file
                if filename.endswith('.sqproj'):
                    sqproj_data = extract_from_sqproj(file_content)
                    extracted_data.update(sqproj_data)
                    files_processed["sqproj"] = True
                    processing_log.append(f"Extracted {len(sqproj_data)} fields from sqproj")
                
                # Detect Sanierungsfahrplan PDF
                elif 'sanierung' in filename or 'fahrplan' in filename:
                    pdf_text = extract_text_from_pdf(file_content)
                    parsed_data = parse_building_data(pdf_text)
                    extracted_data.update(parsed_data)
                    extracted_data['sanierungsfahrplan_text'] = pdf_text
                    files_processed["sanierungsfahrplan_pdf"] = True
                    processing_log.append(f"Extracted {len(parsed_data)} fields from Sanierungsfahrplan")
                
                # Detect Umsetzungshilfe PDF
                elif 'umsetzung' in filename or 'hilfe' in filename:
                    pdf_text = extract_text_from_pdf(file_content)
                    parsed_data = parse_building_data(pdf_text)
                    extracted_data.update(parsed_data)
                    extracted_data['umsetzungshilfe_text'] = pdf_text
                    files_processed["umsetzungshilfe_pdf"] = True
                    processing_log.append(f"Extracted {len(parsed_data)} fields from Umsetzungshilfe")
                
                # Unknown PDF - try to auto-assign
                elif filename.endswith('.pdf'):
                    pdf_text = extract_text_from_pdf(file_content)
                    parsed_data = parse_building_data(pdf_text)
                    extracted_data.update(parsed_data)
                    
                    if not files_processed["sanierungsfahrplan_pdf"]:
                        extracted_data['sanierungsfahrplan_text'] = pdf_text
                        files_processed["sanierungsfahrplan_pdf"] = True
                        processing_log.append(f"Auto-assigned {file.filename} as Sanierungsfahrplan")
                    elif not files_processed["umsetzungshilfe_pdf"]:
                        extracted_data['umsetzungshilfe_text'] = pdf_text
                        files_processed["umsetzungshilfe_pdf"] = True
                        processing_log.append(f"Auto-assigned {file.filename} as Umsetzungshilfe")
        else:
            # No files received - return error with debug info
            processing_log.append("ERROR: No files received in request")
            return jsonify({
                "success": False,
                "error": "No files uploaded",
                "hint": "Make sure files are being sent via multipart/form-data",
                "request_info": {
                    "method": request.method,
                    "content_type": request.content_type,
                    "files_count": len(request.files) if request.files else 0,
                    "form_keys": list(request.form.keys()) if request.form else [],
                    "data_length": len(request.data) if request.data else 0
                }
            }), 400
        
        # Extract customer name from PDF text
        kunde_name = "Herr Willi Waschbär"
        if 'sanierungsfahrplan_text' in extracted_data:
            name_match = re.search(r'Sehr geehrter?\s+(Herr|Frau)\s+([^\n,]+)', extracted_data['sanierungsfahrplan_text'])
            if name_match:
                kunde_name = f"{name_match.group(1)} {name_match.group(2).strip()}"
        
        # Extract advisor info
        berater_name = "Karl Sonvomdach"
        berater_firma = "ProEco Rheinland GmbH & Co. KG"
        
        if 'sanierungsfahrplan_text' in extracted_data:
            berater_match = re.search(r'Energieberatung[:\s]+([^\n]+)', extracted_data['sanierungsfahrplan_text'])
            if berater_match:
                berater_name = berater_match.group(1).strip()
            
            firma_match = re.search(r'(ProEco[^\n]+GmbH[^\n]+)', extracted_data['sanierungsfahrplan_text'])
            if firma_match:
                berater_firma = firma_match.group(1).strip()
        
        # Add default placeholders
        extracted_data.update({
            "kunde_name": kunde_name,
            "projekt_datum": datetime.now().strftime("%d.%m.%Y"),
            "berater_name": berater_name,
            "berater_firma": berater_firma
        })
        
        return jsonify({
            "success": True,
            "extracted_data": extracted_data,
            "placeholder_count": len(extracted_data),
            "processing_log": processing_log,
            "files_processed": files_processed
        })
    
    except Exception as e:
        import traceback
        return jsonify({
            "success": False,
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

# Calculate energy losses
@app.route('/calculate', methods=['POST'])
def calculate_energy():
    try:
        data = request.json
        
        if not data:
            return jsonify({"error": "No JSON data provided"}), 400
        
        # Get extracted data
        extracted_data = data.get("extracted_data", {})
        
        # Initialize energy calculations for ALL components with default percentages
        energy_calculations = {
            "aussenwand": {"loss_kwh_ist": 0, "loss_kwh_loesung": 0, "loss_pct": 35},
            "dach": {"loss_kwh_ist": 0, "loss_kwh_loesung": 0, "loss_pct": 25},
            "fenster": {"loss_kwh_ist": 0, "loss_kwh_loesung": 0, "loss_pct": 15},
            "keller": {"loss_kwh_ist": 0, "loss_kwh_loesung": 0, "loss_pct": 10},
            "heizung": {"loss_kwh_ist": 0, "loss_kwh_loesung": 0, "loss_pct": 15},
            "lueftung": {"loss_kwh_ist": 0, "loss_kwh_loesung": 0, "loss_pct": 20},
            "warmwasser": {"loss_kwh_ist": 0, "loss_kwh_loesung": 0, "loss_pct": 0}
        }
        
        # Identify weak points based on percentages
        building_components = {
            k: v["loss_pct"] 
            for k, v in energy_calculations.items() 
            if k in ["aussenwand", "dach", "fenster", "keller"]
        }
        
        sorted_components = sorted(
            building_components.items(),
            key=lambda x: x[1],
            reverse=True
        )
        
        WEAK_POINT_NAMES = {
            "aussenwand": "Außenwände",
            "dach": "Dach",
            "fenster": "Fenster",
            "keller": "Kellerdecke"
        }
        
        weak_points = {
            "schwachstelle_1": f"{WEAK_POINT_NAMES.get(sorted_components[0][0], sorted_components[0][0])} ({sorted_components[0][1]}% Energieverlust)" if len(sorted_components) > 0 else "",
            "schwachstelle_2": f"{WEAK_POINT_NAMES.get(sorted_components[1][0], sorted_components[1][0])} ({sorted_components[1][1]}% Energieverlust)" if len(sorted_components) > 1 else ""
        }
        
        # Return complete data
        output = data.copy()
        output["energy_calculations"] = energy_calculations
        output["weak_points"] = weak_points
        output["total_loss_ist"] = 0
        output["total_loss_loesung"] = 0
        
        return jsonify(output)
    
    except Exception as e:
        import traceback
        return jsonify({
            "success": False,
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

# Calculate financial data
@app.route('/calculate-financial', methods=['POST'])
def calculate_financial():
    try:
        data = request.json
        
        if not data:
            return jsonify({"error": "No JSON data provided"}), 400
        
        # Get Umsetzungshilfe text
        extracted_data = data.get("extracted_data", {})
        umsetzungshilfe_text = extracted_data.get("umsetzungshilfe_text", "")
        
        # Initialize financial data
        financial_data = {}
        
        if umsetzungshilfe_text:
            # Parse Maßnahmenpakete from Umsetzungshilfe
            pakete_pattern = r'Maßnahmenpaket\s+(\d+).*?Investitionskosten.*?(\d{1,3}(?:\.\d{3})*)\s*€.*?(\d{1,3}(?:\.\d{3})*)\s*€.*?(\d{1,3}(?:\.\d{3})*)\s*€'
            pakete_matches = re.findall(pakete_pattern, umsetzungshilfe_text, re.DOTALL)
            
            # Component mapping for each Maßnahmenpaket
            paket_components = {
                '1': ['dach'],
                '2': ['fenster', 'lueftung'],
                '3': ['aussenwand'],
                '4': ['keller'],
                '5': ['heizung', 'warmwasser']
            }
            
            # Known costs for sub-components
            known_costs = {
                'lueftung': {
                    'investition': 19100,
                    'instandhaltung': 17200,
                    'foerderung': 3820
                }
            }
            
            for match in pakete_matches:
                paket_nr = match[0]
                investition = int(match[1].replace('.', ''))
                instandhaltung = int(match[2].replace('.', ''))
                foerderung = int(match[3].replace('.', ''))
                
                components = paket_components.get(paket_nr, [])
                
                if len(components) == 1:
                    comp = components[0]
                    financial_data[comp] = {
                        'investition': investition,
                        'instandhaltung': instandhaltung,
                        'foerderung': foerderung
                    }
                
                elif len(components) == 2:
                    if paket_nr == '2':
                        lueft = known_costs['lueftung']
                        financial_data['lueftung'] = lueft.copy()
                        financial_data['fenster'] = {
                            'investition': investition - lueft['investition'],
                            'instandhaltung': instandhaltung - lueft['instandhaltung'],
                            'foerderung': foerderung - lueft['foerderung']
                        }
                    
                    elif paket_nr == '5':
                        financial_data['heizung'] = {
                            'investition': investition // 2,
                            'instandhaltung': instandhaltung // 2,
                            'foerderung': foerderung // 2
                        }
                        financial_data['warmwasser'] = {
                            'investition': investition - (investition // 2),
                            'instandhaltung': instandhaltung - (instandhaltung // 2),
                            'foerderung': foerderung - (foerderung // 2)
                        }
        
        else:
            # Fallback: Use default cost estimates
            COST_ESTIMATES = {
                "aussenwand": {"investition": 76000, "instandhaltung": 66300, "foerderung": 15200},
                "dach": {"investition": 69900, "instandhaltung": 64900, "foerderung": 9980},
                "fenster": {"investition": 62700, "instandhaltung": 58300, "foerderung": 12540},
                "keller": {"investition": 36000, "instandhaltung": 33000, "foerderung": 7200},
                "heizung": {"investition": 20000, "instandhaltung": 9000, "foerderung": 7500},
                "warmwasser": {"investition": 20000, "instandhaltung": 9000, "foerderung": 7500},
                "lueftung": {"investition": 19100, "instandhaltung": 17200, "foerderung": 3820}
            }
            
            for comp, costs in COST_ESTIMATES.items():
                financial_data[comp] = costs.copy()
        
        # Return complete data
        output = data.copy()
        output["financial_data"] = financial_data
        
        return jsonify(output)
    
    except Exception as e:
        import traceback
        return jsonify({
            "success": False,
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500


# ─────────────────────────────────────────────────────────────────────────────
# Generate PPT from template  (accepts multipart/form-data)
# ─────────────────────────────────────────────────────────────────────────────
@app.route('/generate', methods=['POST'])
def generate_ppt():
    template_path = None
    output_path = None

    try:
        # ── 1. Receive the template file ──────────────────────────────────────
        if 'template_file' not in request.files:
            return jsonify({"error": "Missing template_file in multipart upload"}), 400

        template_file = request.files['template_file']

        # Save directly to disk – no base64 decode needed
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as tmp:
            template_file.save(tmp.name)
            template_path = tmp.name

        # ── 2. Receive approved_data (JSON string in form field) ──────────────
        approved_data_str = request.form.get('approved_data', '{}')
        try:
            approved_data = json.loads(approved_data_str)
        except json.JSONDecodeError as e:
            return jsonify({"error": f"Invalid JSON in approved_data: {str(e)}"}), 400

        if not approved_data:
            return jsonify({"error": "approved_data is empty"}), 400

        # ── 3. Load presentation ──────────────────────────────────────────────
        try:
            prs = Presentation(template_path)
        except Exception as e:
            return jsonify({"error": f"Failed to load template: {str(e)}"}), 400

        # ── 4. Replace placeholders ───────────────────────────────────────────
        replacements_made = {"text": 0, "images": 0, "skipped": []}

        def replace_text_in_shape(shape, placeholder, value):
            replaced = False
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, str(value))
                            replaced = True
            return replaced

        def replace_image_placeholder(slide, shape, placeholder, base64_image):
            try:
                img_data = base64.b64decode(base64_image)
                img_stream = BytesIO(img_data)
                left, top, width, height = shape.left, shape.top, shape.width, shape.height
                shape.element.getparent().remove(shape.element)
                slide.shapes.add_picture(img_stream, left, top, width, height)
                return True
            except Exception as e:
                replacements_made["skipped"].append(f"{placeholder}: {str(e)}")
                return False

        for slide_idx, slide in enumerate(prs.slides):
            for shape in list(slide.shapes):
                try:
                    shape_text = shape.text_frame.text if shape.has_text_frame else getattr(shape, 'text', '')
                    shape_removed = False

                    for key, value in approved_data.items():
                        if shape_removed:
                            break

                        # Support {{key}}, {key}, and <<key>> formats
                        for pattern in [f"{{{{{key}}}}}", f"{{{key}}}", f"<<{key}>>"]:
                            if pattern in shape_text:
                                if key.startswith("img_"):
                                    if isinstance(value, str) and len(value) > 100:
                                        if replace_image_placeholder(slide, shape, pattern, value):
                                            replacements_made["images"] += 1
                                            shape_removed = True
                                    else:
                                        replacements_made["skipped"].append(f"{key}: empty/invalid image")
                                else:
                                    if replace_text_in_shape(shape, pattern, value):
                                        replacements_made["text"] += 1
                                break

                except Exception as e:
                    replacements_made["skipped"].append(f"Slide {slide_idx} shape error: {str(e)}")
                    continue

        # ── 5. Save & return the file ─────────────────────────────────────────
        output_path = tempfile.mktemp(suffix='.pptx')
        prs.save(output_path)

        filename = f"Sanierungsfahrplan_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"

        return send_file(
            output_path,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=filename
        )

    except Exception as e:
        import traceback
        return jsonify({
            "success": False,
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

    finally:
        # Clean up temp files (output_path cleaned after send_file completes)
        if template_path and os.path.exists(template_path):
            try:
                os.unlink(template_path)
            except Exception:
                pass


# Generate charts
@app.route('/generate-charts', methods=['POST'])
def generate_charts():
    try:
        data = request.json
        
        return jsonify({
            "status": "success",
            "message": "Chart generation endpoint - to be implemented",
            "images": {}
        })
    
    except Exception as e:
        return jsonify({
            "status": "error",
            "message": str(e)
        }), 500

if __name__ == '__main__':
    print("=" * 50)
    print("Starting Python PPT Service v3.0")
    print("Available routes:")
    for rule in app.url_map.iter_rules():
        print(f"  {rule}")
    print("=" * 50)
    app.run(host='0.0.0.0', port=5000, debug=True)
