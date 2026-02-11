"""
Flask App v4.0 - Comprehensive Data Extraction & PPT Generation
Extracts ALL data from .sqproj and PDFs for intelligent PPT placeholder filling
"""

from flask import Flask, request, jsonify, send_file
import zipfile
import xml.etree.ElementTree as ET
import os
import tempfile
import base64
from datetime import datetime
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
import PyPDF2
import json
import re

# Matplotlib for chart generation
import matplotlib
matplotlib.use('Agg')  # Non-interactive backend
import matplotlib.pyplot as plt
import numpy as np

app = Flask(__name__)

# ============================================================================
# ROOT & HEALTH ENDPOINTS
# ============================================================================

@app.route('/')
def home():
    return jsonify({
        "status": "Comprehensive Data Extraction Service v4.0",
        "version": "4.0",
        "description": "Extracts ALL data from building renovation documents",
        "endpoints": [
            "/health",
            "/extract-comprehensive",
            "/read-template-placeholders",
            "/generate-charts",
            "/generate"
        ]
    })

@app.route('/health')
def health():
    return jsonify({
        "status": "healthy",
        "timestamp": datetime.now().isoformat(),
        "version": "4.0"
    })


# ============================================================================
# HELPER: Comprehensive PDF Parser
# ============================================================================

def extract_all_data_from_pdf(pdf_bytes):
    """
    Extract ALL data from PDF:
    - All text (page by page)
    - All numbers with context
    - Metadata
    """
    try:
        pdf_file = BytesIO(pdf_bytes)
        pdf_reader = PyPDF2.PdfReader(pdf_file)
        
        result = {
            'pages': [],
            'full_text': '',
            'metadata': {},
            'numbers_found': []
        }
        
        # Extract metadata
        if pdf_reader.metadata:
            result['metadata'] = {
                'title': pdf_reader.metadata.get('/Title', ''),
                'author': pdf_reader.metadata.get('/Author', ''),
                'creator': pdf_reader.metadata.get('/Creator', ''),
                'producer': pdf_reader.metadata.get('/Producer', ''),
                'creation_date': pdf_reader.metadata.get('/CreationDate', '')
            }
        
        # Extract page-by-page text
        for i, page in enumerate(pdf_reader.pages, 1):
            page_text = page.extract_text()
            result['pages'].append({
                'page_number': i,
                'text': page_text
            })
            result['full_text'] += page_text + '\n\n'
        
        # Extract all numbers with context
        number_pattern = r'([\w\s]{0,30})([\d\.,]+)\s*(€|kWh|m²|kg|%|W|cm|mm|a)?'
        for match in re.finditer(number_pattern, result['full_text']):
            result['numbers_found'].append({
                'context': match.group(1).strip(),
                'value': match.group(2),
                'unit': match.group(3) or ''
            })
        
        return result
        
    except Exception as e:
        return {'error': f'PDF extraction failed: {str(e)}'}


# ============================================================================
# HELPER: Comprehensive .sqproj Parser
# ============================================================================

def extract_all_from_sqproj(sqproj_bytes):
    """
    Extract ALL data from .sqproj file:
    - All XML elements with hierarchy
    - All attributes
    - All text values
    - File structure
    """
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.sqproj') as tmp:
            tmp.write(sqproj_bytes)
            tmp_path = tmp.name
        
        result = {
            'xml_data': {},
            'file_list': [],
            'attributes': {},
            'text_values': []
        }
        
        try:
            with zipfile.ZipFile(tmp_path, 'r') as zip_ref:
                result['file_list'] = zip_ref.namelist()
                
                for filename in zip_ref.namelist():
                    if filename.endswith('.xml'):
                        xml_content = zip_ref.read(filename)
                        try:
                            root = ET.fromstring(xml_content)
                            
                            # Extract all elements with hierarchy
                            def parse_element(elem, path=''):
                                current_path = f'{path}/{elem.tag}' if path else elem.tag
                                
                                # Store attributes
                                if elem.attrib:
                                    result['attributes'][current_path] = elem.attrib
                                
                                # Store text
                                if elem.text and elem.text.strip():
                                    text_val = elem.text.strip()
                                    result['text_values'].append({
                                        'path': current_path,
                                        'value': text_val
                                    })
                                    
                                    # Also store in structured format
                                    key = current_path.split('/')[-1].lower().replace(' ', '_')
                                    result['xml_data'][key] = text_val
                                
                                # Recurse into children
                                for child in elem:
                                    parse_element(child, current_path)
                            
                            parse_element(root)
                            
                        except Exception as e:
                            result['xml_parse_error'] = str(e)
        
        except Exception as e:
            result['sqproj_error'] = str(e)
        
        os.unlink(tmp_path)
        return result
        
    except Exception as e:
        return {'error': f'Failed to process sqproj: {str(e)}'}


# ============================================================================
# EXTRACTION HELPERS: Individual Data Categories
# ============================================================================

def extract_building_info(text):
    """Extract building identification and basic info"""
    info = {}
    
    # Address
    m = re.search(r'([\w\s]{3,40}(?:str|Str|weg|Weg|gasse|Allee)\s*\d+[^\n]*\n?\s*\d{5}\s+\w+)', text)
    if m:
        info['address'] = re.sub(r'\s+', ' ', m.group(1)).strip()
    
    # Owner
    m = re.search(r'Sehr geehrter?\s+(Herr|Frau)\s+([^\n,]+)', text)
    if m:
        info['owner'] = f"{m.group(1)} {m.group(2).strip()}"
    
    # Building type
    m = re.search(r'Gebäudetyp\s+([^\n]+)', text)
    if m:
        info['type'] = m.group(1).strip()
    
    # Construction year
    m = re.search(r'Baujahr\s+(\d{4})', text)
    if m:
        info['construction_year'] = m.group(1)
    
    # Living area
    m = re.search(r'Wohnfläche\s+ca\.\s*([\d,\.]+)\s*m', text)
    if m:
        info['living_area_m2'] = m.group(1).replace(',', '.')
    
    # Floors
    m = re.search(r'Vollgeschosse\s+(\d+)', text)
    if m:
        info['floors'] = m.group(1)
    
    # Basement
    m = re.search(r'Keller\s+(ja[^\n]*|nein)', text, re.IGNORECASE)
    if m:
        info['basement'] = m.group(1).strip()
    
    return info


def extract_consultant_info(text):
    """Extract consultant/advisor information"""
    info = {}
    
    # Name
    m = re.search(r'Energieberatung\s*\n\s*([^\n]+)', text)
    if m:
        info['name'] = m.group(1).strip()
    
    # Company
    m = re.search(r'(ProEco[^\n]+(?:GmbH|KG|AG)[^\n]*)', text)
    if m:
        info['company'] = m.group(1).strip()
    
    # BAFA number
    m = re.search(r'Beraternummer[:\s]+([^\n]+)', text)
    if m:
        info['bafa_number'] = m.group(1).strip()
    
    # Vorgangsnr
    m = re.search(r'Vorgangsnr[^\n]*BAFA[^\n]*([A-Z]+\s*\d+)', text)
    if m:
        info['vorgangsnr'] = m.group(1).strip()
    
    return info


def extract_energy_current(text):
    """Extract current (IST) energy state"""
    current = {}
    
    # Find IST section (before "Ziel" or "Zukunft")
    ist_section = text[:text.find('Ihr Haus in Zukun') if 'Ihr Haus in Zukun' in text else len(text)]
    
    # Primary energy demand
    m = re.search(r'Primärenergiebedarf\s*q\s*p?\s*([\d,]+)\s*kWh/\(m²a\)', ist_section)
    if m:
        current['primary_demand_kwh_m2a'] = m.group(1).replace(',', '.')
    
    # End energy consumption
    m = re.search(r'Endenergieverbrauch\s*([\d\.]+)\s*kWh/a', ist_section)
    if m:
        current['end_consumption_kwh_a'] = m.group(1).replace('.', '')
    
    # Energy costs
    m = re.search(r'Energiekosten[³\d\s]*\s*([\d\.]+)\s*€/a', ist_section)
    if m:
        current['costs_eur_a'] = m.group(1).replace('.', '')
    
    # CO2 emissions
    m = re.search(r'äquivalente\s*CO.{0,4}Emission\s*([\d,]+)\s*kg/\(m²a\)', ist_section)
    if m:
        current['co2_kg_m2a'] = m.group(1).replace(',', '.')
    
    return current


def extract_energy_target(text):
    """Extract target (ZIEL) energy state"""
    target = {}
    
    # Find ZIEL section
    if 'Ihr Haus in Zukun' in text:
        ziel_section = text[text.find('Ihr Haus in Zukun'):]
    else:
        return target
    
    # Primary energy demand
    m = re.search(r'Primärenergiebedarf\s*q\s*p?\s*([\d,]+)\s*kWh/\(m²a\)', ziel_section)
    if m:
        target['primary_demand_kwh_m2a'] = m.group(1).replace(',', '.')
    
    # End energy consumption
    m = re.search(r'Endenergieverbrauch\s*([\d\.]+)\s*kWh/a', ziel_section)
    if m:
        target['end_consumption_kwh_a'] = m.group(1).replace('.', '')
    
    # Energy costs
    m = re.search(r'Energiekosten[³\d\s]*\s*([\d\.]+)\s*€/a', ziel_section)
    if m:
        target['costs_eur_a'] = m.group(1).replace('.', '')
    
    # CO2 emissions
    m = re.search(r'äquivalente\s*CO.{0,4}Emission\s*([\d,]+)\s*kg/\(m²a\)', ziel_section)
    if m:
        target['co2_kg_m2a'] = m.group(1).replace(',', '.')
    
    # Efficiency standard
    m = re.search(r'(EH\s*\d+\s*EE)', text)
    if m:
        target['efficiency_standard'] = m.group(1).replace(' ', ' ')
    
    return target


def extract_all_components(text):
    """Extract all component data (walls, roof, windows, etc.)"""
    components = {}
    
    component_names = [
        'aussenwand', 'dach', 'fenster', 'keller', 
        'heizung', 'warmwasser', 'lueftung'
    ]
    
    for comp in component_names:
        components[comp] = {
            'ist_description': '',
            'loesung_description': '',
            'u_value_ist': '',
            'u_value_target': ''
        }
    
    # IST descriptions
    ist_patterns = {
        'aussenwand': r'(?:Außen)?[Ww]ände?[^\n]{0,20}(?:Massiv|gemauert|ungedämmt)[^\n]{0,60}',
        'dach': r'Dach[^\n]{0,20}(?:Dämmung|nicht\s+einsehbar|ungedämmt)[^\n]{0,60}',
        'fenster': r'Fenster[^\n]{0,20}(?:\d-fach|verglast)[^\n]{0,60}',
        'keller': r'Keller[^\n]{0,20}(?:ungedämmt|teilbeheizt)[^\n]{0,60}',
        'heizung': r'(?:Heizung|Erzeuger)[^\n]{0,20}(?:Gas|Kessel|Standardkessel)[^\n]{0,60}',
        'warmwasser': r'Warmwasser[^\n]{0,20}(?:über|Wärmeerzeuger)[^\n]{0,60}',
        'lueftung': r'Lüftung[^\n]{0,20}(?:Freie|Fenster|Infiltration)[^\n]{0,60}'
    }
    
    for comp, pattern in ist_patterns.items():
        m = re.search(pattern, text, re.IGNORECASE)
        if m:
            components[comp]['ist_description'] = m.group(0).strip()
    
    # LÖSUNG descriptions
    loesung_patterns = {
        'aussenwand': r'Außenwand[^\n]{0,20}Dämmung[^\n]{0,60}',
        'dach': r'Dach[^\n]{0,20}(?:ZSD|ASD|Dämmung)[^\n]{0,60}',
        'fenster': r'Fenster[^\n]{0,20}Uw-Wert[^\n]{0,60}',
        'keller': r'Keller[^\n]{0,20}Dämmung[^\n]{0,60}',
        'heizung': r'(?:Wärmepumpe|Heizung)[^\n]{0,20}(?:Luft-Wasser|L/W)[^\n]{0,60}',
        'warmwasser': r'Warmwasser[^\n]{0,20}(?:Wärmepumpe|Heizungsanlage)[^\n]{0,60}',
        'lueftung': r'Lüftung[^\n]{0,20}(?:WRG|Wärmerückgewinnung)[^\n]{0,60}'
    }
    
    for comp, pattern in loesung_patterns.items():
        m = re.search(pattern, text, re.IGNORECASE)
        if m:
            components[comp]['loesung_description'] = m.group(0).strip()
    
    return components


def extract_measure_packages(text):
    """Extract all 5 measure packages with details"""
    packages = []
    
    # Pattern to find each package
    for i in range(1, 6):
        package = {
            'id': i,
            'name': '',
            'year': '',
            'investment': '',
            'sowieso': '',
            'funding': '',
            'energy_cost_after': ''
        }
        
        # Find package name
        name_pattern = rf'Maßnahmenpaket\s+{i}[^\n]*\n[^\n]*([^\n]{{10,80}})'
        m = re.search(name_pattern, text)
        if m:
            package['name'] = m.group(1).strip()
        
        # Find costs - look for three consecutive euro amounts
        cost_section = re.search(rf'Maßnahmenpaket\s+{i}.*?(\d+\.\d{{3}})\s*€.*?(\d+\.\d{{3}})\s*€.*?(\d+\.\d{{3}})\s*€', text, re.DOTALL)
        if cost_section:
            package['investment'] = cost_section.group(1).replace('.', '')
            package['sowieso'] = cost_section.group(2).replace('.', '')
            package['funding'] = cost_section.group(3).replace('.', '')
        
        # Find year (202X)
        year_match = re.search(rf'Maßnahmenpaket\s+{i}.*?(202\d)', text, re.DOTALL)
        if year_match:
            package['year'] = year_match.group(1)
        
        packages.append(package)
    
    return packages


def extract_costs(text):
    """Extract total costs summary"""
    costs = {}
    
    # Total investment
    m = re.search(r'(?:Gesamtsanierung|gesamt)[^\n]{0,80}(\d+\.\d{3})\s*€', text, re.IGNORECASE)
    if m:
        costs['total_investment'] = m.group(1).replace('.', '')
    
    return costs


def extract_timeline(text):
    """Extract implementation timeline"""
    timeline = {
        'start': '',
        'end': '',
        'packages_schedule': {}
    }
    
    # Look for years in text
    years = re.findall(r'20\d{2}', text)
    if years:
        timeline['start'] = min(years)
        timeline['end'] = max(years)
    
    return timeline


def extract_u_values(text):
    """Extract U-values for all components"""
    u_values = {}
    
    # Pattern: "component ... U-Wert ... number W/(m²K)"
    pattern = r'(\w+)[^\n]{0,40}U-Wert[^\d]{0,10}([\d,\.]+)\s*W/\(m²K\)'
    
    for match in re.finditer(pattern, text, re.IGNORECASE):
        component = match.group(1).lower()
        value = match.group(2).replace(',', '.')
        
        if component not in u_values:
            u_values[component] = []
        u_values[component].append(value)
    
    return u_values


def extract_technical_specs(text):
    """Extract technical specifications"""
    specs = {}
    
    # WLS values
    wls_pattern = r'WLS\s*(\d+)'
    wls_values = re.findall(wls_pattern, text)
    if wls_values:
        specs['wls_values'] = list(set(wls_values))
    
    # JAZ value
    m = re.search(r'JAZ[^\d]{0,10}([\d,\.]+)', text)
    if m:
        specs['jaz'] = m.group(1).replace(',', '.')
    
    # Heat pump type
    m = re.search(r'Wärmepumpe[^\n]{0,40}(Luft[^\n]{0,20}Wasser)', text)
    if m:
        specs['heat_pump_type'] = 'Luft-Wasser'
    
    return specs


# ============================================================================
# HELPER: Structure All Extracted Data
# ============================================================================

def structure_complete_data(sqproj_data, pdf1_data, pdf2_data):
    """
    Merge all extracted data into a unified, intelligently structured format
    ready for AI mapping to PPT placeholders.
    """
    
    # Combine all text for comprehensive search
    all_text = ''
    if pdf1_data.get('full_text'):
        all_text += pdf1_data['full_text'] + '\n\n'
    if pdf2_data.get('full_text'):
        all_text += pdf2_data['full_text'] + '\n\n'
    
    # Normalize whitespace for easier pattern matching
    normalized = re.sub(r'\s+', ' ', all_text)
    
    # Extract key data points with comprehensive patterns
    structured = {
        'raw_data': {
            'sqproj': sqproj_data,
            'pdf1': pdf1_data,
            'pdf2': pdf2_data
        },
        'extracted': {}
    }
    
    # Building Information
    structured['extracted']['building'] = extract_building_info(normalized)
    
    # Consultant Information
    structured['extracted']['consultant'] = extract_consultant_info(normalized)
    
    # Current Energy State
    structured['extracted']['energy_current'] = extract_energy_current(normalized)
    
    # Target Energy State
    structured['extracted']['energy_target'] = extract_energy_target(normalized)
    
    # Components (7 components × IST + LÖSUNG)
    structured['extracted']['components'] = extract_all_components(normalized)
    
    # Measure Packages (5 packages)
    structured['extracted']['measure_packages'] = extract_measure_packages(normalized)
    
    # Costs Summary
    structured['extracted']['costs'] = extract_costs(normalized)
    
    # Timeline
    structured['extracted']['timeline'] = extract_timeline(normalized)
    
    # U-values (all components, before/after)
    structured['extracted']['u_values'] = extract_u_values(normalized)
    
    # Technical specifications
    structured['extracted']['technical_specs'] = extract_technical_specs(normalized)
    
    # All numbers found (for anything we might have missed)
    all_numbers = []
    if pdf1_data.get('numbers_found'):
        all_numbers.extend(pdf1_data['numbers_found'])
    if pdf2_data.get('numbers_found'):
        all_numbers.extend(pdf2_data['numbers_found'])
    structured['extracted']['all_numbers'] = all_numbers[:100]  # Limit to first 100
    
    return structured


# ============================================================================
# HELPER: Read PPT Template Placeholders
# ============================================================================

def read_template_placeholders(template_bytes):
    """
    Read PowerPoint template and extract ALL placeholders
    Returns list of all {{placeholder}} patterns found
    """
    try:
        prs = Presentation(BytesIO(template_bytes))
        placeholders = set()
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    # Find all {{placeholder}} patterns
                    matches = re.findall(r'\{\{([^}]+)\}\}', shape.text)
                    placeholders.update(matches)
                
                # Check tables
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            matches = re.findall(r'\{\{([^}]+)\}\}', cell.text)
                            placeholders.update(matches)
        
        return {
            'placeholders': sorted(list(placeholders)),
            'count': len(placeholders)
        }
        
    except Exception as e:
        return {'error': f'Failed to read template: {str(e)}'}


# ============================================================================
# CHART GENERATION HELPERS
# ============================================================================

def generate_energy_loss_pie_chart(component_name, loss_ist, loss_loesung, loss_pct):
    """Generate a pie chart showing energy loss for a component"""
    fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 4))
    
    # IST state
    colors_ist = ['#dc3545', '#6c757d']
    ax1.pie([loss_ist, 100-loss_ist], labels=['Verlust', 'Effizienz'], 
            colors=colors_ist, autopct='%1.1f%%', startangle=90)
    ax1.set_title(f'{component_name} - IST-Zustand')
    
    # LÖSUNG state
    colors_loesung = ['#28a745', '#6c757d']
    ax2.pie([loss_loesung, 100-loss_loesung], labels=['Verlust', 'Effizienz'], 
            colors=colors_loesung, autopct='%1.1f%%', startangle=90)
    ax2.set_title(f'{component_name} - LÖSUNG')
    
    plt.tight_layout()
    
    # Convert to base64
    buffer = BytesIO()
    plt.savefig(buffer, format='png', dpi=150, bbox_inches='tight')
    buffer.seek(0)
    image_base64 = base64.b64encode(buffer.read()).decode()
    plt.close()
    
    return image_base64


def generate_energy_efficiency_scale(current_value, target_value, max_value=300):
    """Generate energy efficiency scale visualization"""
    fig, ax = plt.subplots(figsize=(12, 3))
    
    # Color gradient from green to red
    colors = ['#28a745', '#5cb85c', '#ffc107', '#fd7e14', '#dc3545']
    ranges = [0, 60, 90, 130, 180, max_value]
    
    # Draw colored bars
    for i in range(len(colors)):
        ax.barh(0, ranges[i+1] - ranges[i], left=ranges[i], 
                height=0.5, color=colors[i], alpha=0.7)
    
    # Mark current and target values
    ax.plot(current_value, 0, 'rv', markersize=15, label=f'IST: {current_value} kWh/(m²a)')
    ax.plot(target_value, 0, 'g^', markersize=15, label=f'ZIEL: {target_value} kWh/(m²a)')
    
    ax.set_xlim(0, max_value)
    ax.set_ylim(-0.5, 0.5)
    ax.set_xlabel('Primärenergiebedarf [kWh/(m²a)]')
    ax.set_yticks([])
    ax.legend(loc='upper right')
    ax.set_title('Energieeffizienz-Skala')
    ax.grid(True, alpha=0.3)
    
    plt.tight_layout()
    
    buffer = BytesIO()
    plt.savefig(buffer, format='png', dpi=150, bbox_inches='tight')
    buffer.seek(0)
    image_base64 = base64.b64encode(buffer.read()).decode()
    plt.close()
    
    return image_base64


def generate_timeline_chart(measure_packages):
    """Generate Gantt-style timeline chart for measure packages"""
    if not measure_packages or len(measure_packages) == 0:
        return None
    
    fig, ax = plt.subplots(figsize=(12, 6))
    
    # Extract data
    packages = []
    years = []
    for pkg in measure_packages:
        if pkg.get('name') and pkg.get('year'):
            packages.append(pkg['name'])
            try:
                years.append(int(pkg['year']))
            except:
                years.append(2025)
    
    if not packages:
        return None
    
    # Create timeline
    min_year = min(years) if years else 2025
    max_year = max(years) if years else 2029
    
    colors = ['#007bff', '#28a745', '#ffc107', '#dc3545', '#6f42c1']
    
    for i, (pkg, year) in enumerate(zip(packages, years)):
        ax.barh(i, 1, left=year-min_year, height=0.5, 
                color=colors[i % len(colors)], alpha=0.8)
        ax.text(year-min_year+0.5, i, pkg, va='center', fontsize=9)
    
    ax.set_yticks(range(len(packages)))
    ax.set_yticklabels([f'MP{i+1}' for i in range(len(packages))])
    ax.set_xticks(range(max_year - min_year + 1))
    ax.set_xticklabels([str(min_year + i) for i in range(max_year - min_year + 1)])
    ax.set_xlabel('Jahr')
    ax.set_title('Sanierungsfahrplan - Zeitliche Umsetzung')
    ax.grid(True, alpha=0.3, axis='x')
    
    plt.tight_layout()
    
    buffer = BytesIO()
    plt.savefig(buffer, format='png', dpi=150, bbox_inches='tight')
    buffer.seek(0)
    image_base64 = base64.b64encode(buffer.read()).decode()
    plt.close()
    
    return image_base64


def generate_cost_comparison_chart(measure_packages):
    """Generate cost comparison bar chart"""
    if not measure_packages or len(measure_packages) == 0:
        return None
    
    fig, ax = plt.subplots(figsize=(12, 6))
    
    # Extract data
    labels = []
    investments = []
    sowiesos = []
    fundings = []
    
    for i, pkg in enumerate(measure_packages):
        labels.append(f'MP{i+1}')
        try:
            inv = pkg.get('investment', '0').replace('.', '').replace(',', '.')
            investments.append(float(inv) if inv else 0)
            
            sow = pkg.get('sowieso', '0').replace('.', '').replace(',', '.')
            sowiesos.append(float(sow) if sow else 0)
            
            fun = pkg.get('funding', '0').replace('.', '').replace(',', '.')
            fundings.append(float(fun) if fun else 0)
        except:
            investments.append(0)
            sowiesos.append(0)
            fundings.append(0)
    
    if sum(investments) == 0:
        return None
    
    x = np.arange(len(labels))
    width = 0.25
    
    ax.bar(x - width, investments, width, label='Investition', color='#dc3545')
    ax.bar(x, sowiesos, width, label='Sowieso-Kosten', color='#ffc107')
    ax.bar(x + width, fundings, width, label='Förderung', color='#28a745')
    
    ax.set_xlabel('Maßnahmenpaket')
    ax.set_ylabel('Kosten [€]')
    ax.set_title('Kostenvergleich der Maßnahmenpakete')
    ax.set_xticks(x)
    ax.set_xticklabels(labels)
    ax.legend()
    ax.grid(True, alpha=0.3, axis='y')
    
    # Format y-axis as currency
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda x, p: f'{int(x):,} €'.replace(',', '.')))
    
    plt.tight_layout()
    
    buffer = BytesIO()
    plt.savefig(buffer, format='png', dpi=150, bbox_inches='tight')
    buffer.seek(0)
    image_base64 = base64.b64encode(buffer.read()).decode()
    plt.close()
    
    return image_base64


# ============================================================================
# API ENDPOINTS
# ============================================================================

@app.route('/extract-comprehensive', methods=['POST'])
def extract_comprehensive():
    """
    Extract ALL data from .sqproj and 2 PDFs
    Returns unified structured data ready for intelligent mapping
    """
    try:
        if 'files' not in request.files:
            return jsonify({'error': 'No files uploaded'}), 400
        
        files = request.files.getlist('files')
        
        if len(files) < 3:
            return jsonify({'error': 'Expected 3 files (2 PDFs + 1 .sqproj)'}), 400
        
        # Identify files by type
        pdf_files = []
        sqproj_file = None
        
        for file in files:
            if file.filename.endswith('.pdf'):
                pdf_files.append(file.read())
            elif file.filename.endswith('.sqproj'):
                sqproj_file = file.read()
        
        if not sqproj_file or len(pdf_files) < 2:
            return jsonify({'error': 'Expected 2 PDFs and 1 .sqproj file'}), 400
        
        # Extract from each file
        sqproj_data = extract_all_from_sqproj(sqproj_file)
        pdf1_data = extract_all_data_from_pdf(pdf_files[0])
        pdf2_data = extract_all_data_from_pdf(pdf_files[1])
        
        # Structure all data
        complete_data = structure_complete_data(sqproj_data, pdf1_data, pdf2_data)
        
        return jsonify({
            'success': True,
            'data': complete_data,
            'extraction_summary': {
                'sqproj_files_found': len(sqproj_data.get('file_list', [])),
                'pdf1_pages': len(pdf1_data.get('pages', [])),
                'pdf2_pages': len(pdf2_data.get('pages', [])),
                'total_data_points': len(str(complete_data))
            }
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/read-template-placeholders', methods=['POST'])
def api_read_template():
    """
    Read PowerPoint template and extract all placeholders
    """
    try:
        if 'template' not in request.files:
            return jsonify({'error': 'No template file uploaded'}), 400
        
        template_file = request.files['template']
        template_bytes = template_file.read()
        
        result = read_template_placeholders(template_bytes)
        
        return jsonify({
            'success': True,
            'result': result
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/generate-charts', methods=['POST'])
def generate_charts():
    """
    Generate all chart images for PowerPoint placeholders
    Returns base64-encoded PNG images
    """
    try:
        data = request.get_json()
        
        if not data:
            return jsonify({'error': 'No data provided'}), 400
        
        charts = {}
        
        # Extract energy data
        energy_current = data.get('energy_current', {})
        energy_target = data.get('energy_target', {})
        components = data.get('components', {})
        measure_packages = data.get('measure_packages', [])
        
        # Generate energy efficiency scale
        try:
            current_primary = float(energy_current.get('primary_demand_kwh_m2a', 290))
            target_primary = float(energy_target.get('primary_demand_kwh_m2a', 27))
            charts['img_energieklasse_ist'] = generate_energy_efficiency_scale(
                current_primary, current_primary, 300
            )
            charts['img_energieklasse_ziel'] = generate_energy_efficiency_scale(
                target_primary, target_primary, 300
            )
        except Exception as e:
            print(f"Error generating efficiency scale: {e}")
        
        # Generate component energy loss charts
        component_mapping = {
            'aussenwand': 'img_energieverluste_aussenwand',
            'dach': 'img_energieverluste_dach',
            'fenster': 'img_energieverluste_fenster',
            'keller': 'img_energieverluste_keller',
            'heizung': 'img_energieverluste_heizung'
        }
        
        for comp_key, img_key in component_mapping.items():
            try:
                # Mock data - you would calculate actual losses
                loss_ist = 30  # percentage
                loss_loesung = 5  # percentage
                loss_pct = 25
                
                chart = generate_energy_loss_pie_chart(
                    comp_key.capitalize(), 
                    loss_ist, 
                    loss_loesung, 
                    loss_pct
                )
                charts[img_key] = chart
            except Exception as e:
                print(f"Error generating chart for {comp_key}: {e}")
        
        # Generate timeline chart
        try:
            timeline_chart = generate_timeline_chart(measure_packages)
            if timeline_chart:
                charts['img_sanierungsfahrplan'] = timeline_chart
        except Exception as e:
            print(f"Error generating timeline: {e}")
        
        # Generate cost comparison chart
        try:
            cost_chart = generate_cost_comparison_chart(measure_packages)
            if cost_chart:
                charts['img_kostenvergleich'] = cost_chart
        except Exception as e:
            print(f"Error generating cost chart: {e}")
        
        return jsonify({
            'success': True,
            'images': charts,
            'count': len(charts)
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/generate', methods=['POST'])
def generate_ppt():
    """
    Generate PowerPoint by filling template with mapped data
    """
    try:
        data = request.get_json()
        
        if not data:
            return jsonify({'error': 'No data provided'}), 400
        
        template_base64 = data.get('template_file')
        approved_data = data.get('approved_data')
        
        if not template_base64 or not approved_data:
            return jsonify({'error': 'Missing template_file or approved_data'}), 400
        
        # Decode template
        template_bytes = base64.b64decode(template_base64)
        prs = Presentation(BytesIO(template_bytes))
        
        # Track replacements
        replacements = {}
        
        # Replace text placeholders
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    original_text = shape.text
                    new_text = original_text
                    
                    for placeholder, value in approved_data.items():
                        pattern = f'{{{{{placeholder}}}}}'
                        if pattern in new_text:
                            # Skip image placeholders in text replacement
                            if not placeholder.startswith('img_'):
                                new_text = new_text.replace(pattern, str(value))
                                replacements[placeholder] = replacements.get(placeholder, 0) + 1
                    
                    if hasattr(shape, 'text_frame'):
                        shape.text_frame.text = new_text
                
                # Handle tables
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            original_text = cell.text
                            new_text = original_text
                            
                            for placeholder, value in approved_data.items():
                                pattern = f'{{{{{placeholder}}}}}'
                                if pattern in new_text:
                                    if not placeholder.startswith('img_'):
                                        new_text = new_text.replace(pattern, str(value))
                                        replacements[placeholder] = replacements.get(placeholder, 0) + 1
                            
                            cell.text = new_text
                
                # Handle image placeholders
                if hasattr(shape, 'text') and '{{img_' in shape.text:
                    for placeholder, value in approved_data.items():
                        if placeholder.startswith('img_') and f'{{{{{placeholder}}}}}' in shape.text:
                            try:
                                # Decode base64 image
                                img_data = base64.b64decode(value)
                                img_stream = BytesIO(img_data)
                                
                                # Replace with actual image
                                left = shape.left
                                top = shape.top
                                width = shape.width
                                height = shape.height
                                
                                # Remove the placeholder shape
                                sp = shape.element
                                sp.getparent().remove(sp)
                                
                                # Add the image
                                slide.shapes.add_picture(img_stream, left, top, width, height)
                                replacements[placeholder] = 1
                            except Exception as e:
                                print(f"Error replacing image {placeholder}: {e}")
        
        # Save to BytesIO
        output = BytesIO()
        prs.save(output)
        output.seek(0)
        
        # Encode to base64
        file_content = base64.b64encode(output.read()).decode()
        file_size_mb = len(file_content) / 1024 / 1024
        
        return jsonify({
            'success': True,
            'filename': 'presentation_filled.pptx',
            'file_content': file_content,
            'file_size_mb': round(file_size_mb, 2),
            'replacements': replacements
        })
        
    except Exception as e:
        return jsonify({'error': str(e)}), 500


# ============================================================================
# RUN APP
# ============================================================================

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=8080, debug=True)
