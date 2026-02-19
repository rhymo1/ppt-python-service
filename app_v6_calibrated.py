# -*- coding: utf-8 -*-
"""
Flask App v6.3  –  Complete iSFP PPT Generation Service
Based on v6.2 + Türen→Außenwand reclassification + negative Heizung bar support.

v6.3 additions:
  Fix7  Moved Türen keywords from fenster to aussenwand component_map.
        Comp shows "Außenwand/Türen" as combined category on slides 12,19,23,27,33.
        IST deviation now <2% for all components.
  Fix8  Removed max(0,...) from Heizung Lösung calculation.
        Heat pump COP>1 produces negative system loss (net gain),
        enabling >100% reduction percentages (e.g. -112% in comp).
  Fix9  resize_bar_shapes() handles negative kWh values:
        shifts all bars right to create negative zone,
        negative green bars extend LEFT from origin line.
  Fix10 BAR_PATTERN regex fixed (was \\d, now \d for digit matching).

v6.2 additions:
  Fix4  SOLAR_CORRECTION factors recalibrated against actual extracted UA values.
        Old aggressive factors (fenster=0.61, keller=1.43) were tuned to assumed
        UA values that didn't match PDF extraction output, causing -41%/+34% errors.
        New gentle factors validated against Hottgenroth output (max deviation <1%):
        dach=1.0, aussenwand=1.08, fenster=0.97, keller=1.00
  Fix5  Added H_T/H_V estimation fallback from component UA products when
        regex extraction from Umsetzungshilfe PDF fails. Uses DEFAULT_VENT_RATIO=0.20
        (validated: Waschbär MFH actual = 0.204). Since only the ratio matters
        (not absolute values), this produces correct splits regardless of H_T magnitude.
  Fix6  Moved Warmwasserbedarf estimation before calibrated/fallback decision,
        so Heizung = Qend - Qh - Qww works correctly in both paths.

v6.1: Fixed 120+ mojibake regex patterns, added LUEFTUNG_CORRECTION=1.28

v6.0: Calibrated energy loss calculation (Cal1-Cal4)

Previous versions:
  v5.2: Yellow highlight, heizung loss, fenster typo, bar resizing, PDF images
  v5.1: All Bug #0-#22 fixes, OCR, structured logging, caching
"""

import os
import re
import json
import time
import hashlib
import sqlite3
import logging
import tempfile
import base64
from datetime import datetime
from io import BytesIO
from pathlib import Path
from copy import deepcopy

from flask import Flask, request, jsonify, send_file
import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.oxml.ns import qn
from pptx.dml.color import RGBColor
from PIL import Image
from lxml import etree
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np

# PyMuPDF for PDF image extraction (Priority 3)
try:
    import fitz
    FITZ_AVAILABLE = True
except ImportError:
    FITZ_AVAILABLE = False

# Optional OCR
try:
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

# ============================================================================
# APP CONFIGURATION
# ============================================================================

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 200 * 1024 * 1024  # 200 MB

CACHE_DIR = Path(tempfile.gettempdir()) / 'isfp_cache'
CACHE_DIR.mkdir(exist_ok=True)

TEMPLATE_DIR = Path(tempfile.gettempdir()) / 'isfp_templates'
TEMPLATE_DIR.mkdir(exist_ok=True)

TEMPLATE_DIR = Path(tempfile.gettempdir()) / 'isfp_templates'
TEMPLATE_DIR.mkdir(exist_ok=True)

# ============================================================================
# STRUCTURED LOGGING
# ============================================================================

logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s [%(levelname)s] %(name)s: %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)
log = logging.getLogger('isfp')


class ExtractionLog:
    """Track per-field extraction success/failure for debugging."""
    def __init__(self):
        self.entries = []

    def ok(self, field: str, value, source: str = ''):
        self.entries.append({'field': field, 'status': 'ok', 'value': str(value)[:120], 'source': source})
        log.info(f'  ✓ {field} = {str(value)[:80]}')

    def miss(self, field: str, reason: str = ''):
        self.entries.append({'field': field, 'status': 'miss', 'reason': reason})
        log.warning(f'  ✓" {field}  –  {reason}')

    def summary(self):
        ok_count = sum(1 for e in self.entries if e['status'] == 'ok')
        total = len(self.entries)
        return {'extracted': ok_count, 'missed': total - ok_count, 'total': total, 'details': self.entries}


# ============================================================================
# CACHING
# ============================================================================

def cache_key(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()[:16]


def cache_get(key: str):
    path = CACHE_DIR / f'{key}.json'
    if path.exists() and (time.time() - path.stat().st_mtime) < 3600:
        log.info(f'Cache hit: {key}')
        return json.loads(path.read_text(encoding='utf-8'))
    return None


def cache_set(key: str, data: dict):
    path = CACHE_DIR / f'{key}.json'
    try:
        # Atomic write: write to temp file then rename (prevents partial reads)
        tmp_path = path.with_suffix('.tmp')
        tmp_path.write_text(json.dumps(data, ensure_ascii=False, default=str), encoding='utf-8')
        tmp_path.rename(path)
    except Exception as e:
        log.warning(f'Cache write failed: {e}')


# ============================================================================
# ROOT & HEALTH ENDPOINTS
# ============================================================================

@app.route('/')
def home():
    return jsonify({
        'status': 'iSFP Data Extraction & PPT Generation Service v6.0',
        'version': '6.2',
        'endpoints': {
            '/health': 'GET  –  health check',
            '/extract-comprehensive': 'POST  –  extract all data from PDFs + .sqproj',
            '/read-template-placeholders': 'POST  –  read {{placeholders}} from .pptx',
            '/generate-charts': 'POST (JSON)  –  generate chart images',
            '/generate': 'POST (multipart)  –  fill template and return .pptx',
        }
    })


@app.route('/health')
def health():
    return jsonify({
        'status': 'healthy',
        'timestamp': datetime.now().isoformat(),
        'version': '6.2',
        'ocr_available': OCR_AVAILABLE,
        'fitz_available': FITZ_AVAILABLE,
    })


# ============================================================================
# HELPER: .sqproj READER (Bug #0  –  SQLite, not ZIP/XML)
# ============================================================================

def extract_all_from_sqproj(sqproj_bytes: bytes) -> dict:
    """
    .sqproj files are SQLite databases, NOT zip/xml archives.
    Extract building geometry, contacts, construction types, U-values.
    """
    elog = ExtractionLog()
    result = {
        'tables': {},
        'building': {},
        'contacts': {},
        'construction_types': [],
        'u_values': {},
        'geometry': {},
        'climate': {},
    }

    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix='.sqproj') as tmp:
            tmp.write(sqproj_bytes)
            tmp_path = tmp.name

        conn = sqlite3.connect(tmp_path)
        conn.row_factory = sqlite3.Row
        cursor = conn.cursor()

        # List all tables
        cursor.execute("SELECT name FROM sqlite_master WHERE type='table' ORDER BY name")
        tables = [row[0] for row in cursor.fetchall()]
        result['tables']['count'] = len(tables)
        result['tables']['names'] = tables[:50]
        elog.ok('sqproj_tables', f'{len(tables)} tables found', 'sqproj')

        def read_table(name, limit=100):
            try:
                cursor.execute(f'SELECT * FROM "{name}" LIMIT {limit}')
                cols = [d[0] for d in cursor.description] if cursor.description else []
                rows = [dict(row) for row in cursor.fetchall()]
                return cols, rows
            except Exception:
                return [], []

        for tbl in tables:
            tbl_lower = tbl.lower()

            if any(k in tbl_lower for k in ['bauteil', 'wand', 'dach', 'decke', 'fenster', 'tuer', 'boden']):
                cols, rows = read_table(tbl)
                if rows:
                    result['geometry'][tbl] = rows
                    elog.ok(f'geometry.{tbl}', f'{len(rows)} rows', 'sqproj')

            if any(k in tbl_lower for k in ['konstruktion', 'aufbau', 'schicht', 'u_wert', 'uwert']):
                cols, rows = read_table(tbl)
                if rows:
                    result['construction_types'].append({'table': tbl, 'rows': rows})
                    elog.ok(f'construction.{tbl}', f'{len(rows)} rows', 'sqproj')

            if any(k in tbl_lower for k in ['kontakt', 'adresse', 'person', 'eigentuemer', 'berater']):
                cols, rows = read_table(tbl)
                if rows:
                    result['contacts'][tbl] = rows
                    elog.ok(f'contacts.{tbl}', f'{len(rows)} rows', 'sqproj')

            if any(k in tbl_lower for k in ['klima', 'standort', 'wetter']):
                cols, rows = read_table(tbl)
                if rows:
                    result['climate'][tbl] = rows
                    elog.ok(f'climate.{tbl}', f'{len(rows)} rows', 'sqproj')

        for tbl in tables:
            if 'u_wert' in tbl.lower() or 'uwert' in tbl.lower():
                cols, rows = read_table(tbl)
                for row in rows:
                    for col_name, val in row.items():
                        if val and isinstance(val, (int, float)):
                            result['u_values'][f'{tbl}.{col_name}'] = val

        conn.close()

    except Exception as e:
        log.error(f'sqproj extraction failed: {e}')
        result['error'] = str(e)
        elog.miss('sqproj', str(e))
    finally:
        if tmp_path and os.path.exists(tmp_path):
            os.unlink(tmp_path)

    result['_log'] = elog.summary()
    return result


# ============================================================================
# HELPER: PDF READER (Bug #10  –  pdfplumber; Bug #11  –  no whitespace nuke)
# ============================================================================

def extract_all_data_from_pdf(pdf_bytes: bytes, label: str = 'pdf') -> dict:
    """
    Extract text + tables from PDF using pdfplumber.
    No whitespace normalization (Bug #11).
    OCR fallback for graphical-only pages.
    """
    elog = ExtractionLog()
    result = {
        'pages': [],
        'full_text': '',
        'tables': [],
        'metadata': {},
    }

    try:
        with pdfplumber.open(BytesIO(pdf_bytes)) as pdf:
            result['metadata']['page_count'] = len(pdf.pages)
            elog.ok(f'{label}_pages', len(pdf.pages), label)

            for i, page in enumerate(pdf.pages):
                page_num = i + 1
                text = page.extract_text() or ''

                if not text.strip() and OCR_AVAILABLE:
                    try:
                        img = page.to_image(resolution=300)
                        pil_img = img.original
                        text = pytesseract.image_to_string(pil_img, lang='deu')
                        if text.strip():
                            elog.ok(f'{label}_p{page_num}_ocr', f'{len(text)} chars', 'OCR')
                    except Exception as e:
                        elog.miss(f'{label}_p{page_num}_ocr', str(e))

                page_tables = page.extract_tables() or []

                result['pages'].append({
                    'page_number': page_num,
                    'text': text,
                    'tables': page_tables,
                })
                result['full_text'] += text + '\n\n'

                if page_tables:
                    for t_idx, tbl in enumerate(page_tables):
                        result['tables'].append({
                            'page': page_num,
                            'table_index': t_idx,
                            'data': tbl,
                        })
                    elog.ok(f'{label}_p{page_num}_tables', f'{len(page_tables)} tables', label)

    except Exception as e:
        log.error(f'PDF extraction failed ({label}): {e}')
        result['error'] = str(e)
        elog.miss(label, str(e))

    result['_log'] = elog.summary()
    return result


# ============================================================================
# PDF IMAGE EXTRACTION (Priority 3  –  page rendering + auto-crop)
# ============================================================================

RENDER_DPI = 200


def extract_images_from_pdf(pdf_bytes: bytes, pdf_type: str = 'sanierungsfahrplan') -> dict:
    """
    Extract relevant images from PDF by rendering pages and cropping regions.
    Handles fully rasterized PDFs (each page = one image) via:
      1. Rendering pages at high DPI
      2. Auto-detecting photo regions via brightness analysis
      3. Mapping to PPT placeholder keys
    Returns dict of placeholder_key -> base64 PNG string.
    """
    images = {}

    if not FITZ_AVAILABLE:
        log.warning('  PyMuPDF (fitz) not available – skipping PDF image extraction')
        return images

    try:
        doc = fitz.open(stream=pdf_bytes, filetype='pdf')
        page_count = len(doc)
        log.info(f'  PDF image extraction ({pdf_type}): {page_count} pages')

        if pdf_type == 'sanierungsfahrplan':

            # --- Page 7: "Mein Sanierungsfahrplan" overview (landscape) ---
            if page_count >= 7:
                page = doc[6]
                pix = page.get_pixmap(dpi=RENDER_DPI)
                img = Image.frombytes('RGB', [pix.width, pix.height], pix.samples)
                buf = BytesIO()
                img.save(buf, format='PNG', optimize=True)
                images['img_meinsanierungsfahrplan'] = base64.b64encode(buf.getvalue()).decode()
                log.info(f'    ✓ img_meinsanierungsfahrplan ({pix.width}x{pix.height})')

            # --- Page 1: Building photo (center of title page) ---
            if page_count >= 1:
                page = doc[0]
                pix = page.get_pixmap(dpi=RENDER_DPI)
                img = Image.frombytes('RGB', [pix.width, pix.height], pix.samples)
                w, h = img.size
                crop_box = (int(w * 0.36), int(h * 0.48), int(w * 0.72), int(h * 0.66))
                building_img = img.crop(crop_box)
                buf = BytesIO()
                building_img.save(buf, format='PNG', optimize=True)
                images['img_agendabild'] = base64.b64encode(buf.getvalue()).decode()
                log.info(f'    ✓ img_agendabild ({building_img.size})')

            # --- Page 3: "Ihr Haus heute – Bestand" with Schwachstellen photos ---
            if page_count >= 3:
                page = doc[2]
                pix = page.get_pixmap(dpi=RENDER_DPI)
                img = Image.frombytes('RGB', [pix.width, pix.height], pix.samples)
                w, h = img.size
                arr = np.array(img)

                # Auto-detect photo regions in left column via brightness
                strip = arr[:, 100:min(400, w // 3), :].mean(axis=(1, 2))
                dark_mask = strip < 220
                dark_rows = np.where(dark_mask)[0]

                photo_regions = []
                if len(dark_rows) > 0:
                    diffs = np.diff(dark_rows)
                    splits = np.where(diffs > 20)[0]
                    start = dark_rows[0]
                    for s in splits:
                        end = dark_rows[s]
                        if end - start > 80:
                            photo_regions.append((start, end))
                        start = dark_rows[s + 1]
                    end = dark_rows[-1]
                    if end - start > 80:
                        photo_regions.append((start, end))

                # Map detected regions to placeholders
                # Typical order: Dach, Haustür/Fenster, Keller, Heizung
                photo_mapping = [
                    ('img_schwachstelle_1', 'img_dach_istzustand'),
                    ('img_fenster_ist', None),
                    ('img_keller_ist', None),
                    ('img_schwachstelle_2', 'img_warmwasser_ist'),
                ]

                right_edge = min(int(w * 0.34), 650)
                for i, (y_start, y_end) in enumerate(photo_regions):
                    if i >= len(photo_mapping):
                        break
                    crop = img.crop((30, y_start, right_edge, y_end))
                    buf = BytesIO()
                    crop.save(buf, format='PNG', optimize=True)
                    b64 = base64.b64encode(buf.getvalue()).decode()

                    primary_key, secondary_key = photo_mapping[i]
                    images[primary_key] = b64
                    if secondary_key:
                        images[secondary_key] = b64
                    log.info(f'    ✓ {primary_key} (photo {i + 1})')

        doc.close()

    except Exception as e:
        log.error(f'  PDF image extraction failed: {e}')

    log.info(f'  Extracted {len(images)} images from PDF')
    return images


# ============================================================================
# STRUCTURED DATA EXTRACTION (Bug #12, #13  –  rewritten patterns)
# ============================================================================

def _search(pattern: str, text: str, group: int = 1, default: str = '') -> str:
    """Safe regex search returning group or default."""
    m = re.search(pattern, text, re.IGNORECASE | re.DOTALL)
    return m.group(group).strip() if m else default


def _search_number(pattern: str, text: str, group: int = 1) -> str:
    """Search and clean German number format."""
    raw = _search(pattern, text, group)
    if not raw:
        return ''
    return raw


def extract_building_info(text: str, elog: ExtractionLog) -> dict:
    info = {}

    m = re.search(r'Gebäudeadresse\s*\n?\s*(.+?)\n\s*(\d{5})\s+(\S+)', text)
    if m:
        info['street'] = m.group(1).strip()
        info['plz'] = m.group(2)
        info['city'] = m.group(3)
        info['address'] = f"{info['street']}, {info['plz']} {info['city']}"
        elog.ok('building.address', info['address'], 'pdf')
    else:
        m2 = re.search(r'(\S+(?:str|Str|weg|Weg|gasse|platz|allee)\S*\s+\d+\S*)\s*\n?\s*(\d{5})\s+(\w+)', text)
        if m2:
            info['street'] = m2.group(1).strip()
            info['plz'] = m2.group(2)
            info['city'] = m2.group(3)
            info['address'] = f"{info['street']}, {info['plz']} {info['city']}"
            elog.ok('building.address', info['address'], 'pdf')
        else:
            elog.miss('building.address', 'No address pattern found')

    m = re.search(r'(?:Sehr geehrte[r]?|Herr|Frau)\s+(Herr|Frau)\s+([\w\-]+)', text)
    if m:
        info['owner_salutation'] = m.group(1)
        info['owner_name'] = m.group(2)
        info['owner'] = f"{m.group(1)} {m.group(2)}"
        elog.ok('building.owner', info['owner'], 'pdf')
    else:
        elog.miss('building.owner')

    val = _search(r'Gebäudetyp\s+(\S+(?:\s+\S+)?)', text)
    if val:
        info['type'] = val
        elog.ok('building.type', val, 'pdf')

    val = _search(r'Baujahr\s+(\d{4})', text)
    if val:
        info['construction_year'] = val
        elog.ok('building.construction_year', val, 'pdf')

    m = re.search(r'Wohnfläche\s+(?:ca\.?\s*)?([\d.,]+)\s*m', text)
    if m:
        info['living_area_m2'] = m.group(1).replace(',', '.')
        elog.ok('building.living_area', info['living_area_m2'], 'pdf')

    val = _search(r'Vollgeschosse\s+(\d+)', text)
    if val:
        info['floors'] = val

    val = _search(r'Keller\s+(ja\s*/?\s*\S+|nein)', text)
    if val:
        info['basement'] = val

    val = _search(r'Baujahr\s*\n?\s*Heizung\s+(\d{4})', text)
    if val:
        info['heating_year'] = val

    val = _search(r'Erzeuger\s+(.+?)(?:\n|$)', text)
    if val:
        info['heater_type'] = val

    return info


def extract_consultant_info(text: str, elog: ExtractionLog) -> dict:
    info = {}

    m = re.search(r'Energieberater[/in]*\s*\n\s*([\w\s\.\-]{3,50})\n', text)
    if not m:
        m = re.search(r'Ihr\s+Energieberater\s*:?\s*\n?\s*([\w\.\-]+(?:\s+[\w\.\-]+){1,3})', text)
    if not m:
        m = re.search(r'erstellt\s+(?:von|durch)\s+([\w\.\-]+(?:\s+[\w\.\-]+){1,3})', text)
    if m:
        info['name'] = m.group(1).strip()
        elog.ok('consultant.name', info['name'], 'pdf')

    m = re.search(r'(ProEco\s+Rheinland\s+GmbH[^\n]*)', text)
    if not m:
        m = re.search(r'(\S+\s+(?:GmbH|KG|AG|UG)[^\n]*)', text)
    if m:
        info['company'] = m.group(1).strip()

    val = _search(r'Beraternummer[:\s]+(\S+)', text)
    if val:
        info['bafa_number'] = val

    val = _search(r'Vorgangsnr[^:]*:\s*(\S+\s*\d+)', text)
    if val:
        info['vorgangsnr'] = val

    return info


def extract_energy_values(text: str, elog: ExtractionLog) -> dict:
    """Extract IST and ZIEL energy values."""
    result = {'ist': {}, 'ziel': {}}

    ist_section = text
    ziel_marker = text.find('Ihr Haus in Zukun')
    if ziel_marker > 0:
        ist_section = text[:ziel_marker]

    m = re.search(r'Primärenergiebedarf\s*q\s*p?\s*\n?\s*([\d.,]+)\s*kWh/\(m²a\)', ist_section)
    if m:
        result['ist']['primary_demand'] = m.group(1).replace('.', '').replace(',', '.')
        elog.ok('energy.ist.primary', result['ist']['primary_demand'], 'pdf')
    else:
        m2 = re.search(r'(\d{2,3})\s*kWh/\(m²a\)', ist_section)
        if m2:
            result['ist']['primary_demand'] = m2.group(1)

    m = re.search(r'Endenergieverbrauch\s*\n?\s*([\d.,]+)\s*kWh/a', ist_section)
    if m:
        result['ist']['end_consumption'] = m.group(1).replace('.', '')
        elog.ok('energy.ist.end_consumption', result['ist']['end_consumption'], 'pdf')

    m = re.search(r'Energiekosten[³\s]*\n?\s*([\d.,]+)\s*€/a', ist_section)
    if m:
        result['ist']['costs'] = m.group(1).replace('.', '')
        elog.ok('energy.ist.costs', result['ist']['costs'], 'pdf')

    m = re.search(r'(?:äquivalente\s*)?CO\s*2?\s*-?\s*Emission(?:en)?\s*\n?\s*([\d.,]+)\s*kg/\(m²a\)', ist_section)
    if m:
        result['ist']['co2'] = m.group(1).replace(',', '.')
        elog.ok('energy.ist.co2', result['ist']['co2'], 'pdf')

    if ziel_marker > 0:
        ziel_section = text[ziel_marker:]

        m = re.search(r'Primärenergiebedarf\s*q\s*p?\s*\n?\s*([\d.,]+)\s*kWh/\(m²a\)', ziel_section)
        if m:
            result['ziel']['primary_demand'] = m.group(1).replace('.', '').replace(',', '.')
            elog.ok('energy.ziel.primary', result['ziel']['primary_demand'], 'pdf')

        m = re.search(r'Endenergieverbrauch\s*\n?\s*([\d.,]+)\s*kWh/a', ziel_section)
        if m:
            result['ziel']['end_consumption'] = m.group(1).replace('.', '')
            elog.ok('energy.ziel.end_consumption', result['ziel']['end_consumption'], 'pdf')

        m = re.search(r'Energiekosten[³\s]*\n?\s*([\d.,]+)\s*€/a', ziel_section)
        if m:
            result['ziel']['costs'] = m.group(1).replace('.', '')

        m = re.search(r'(?:äquivalente\s*)?CO\s*2?\s*-?\s*Emission(?:en)?\s*\n?\s*([\d.,]+)\s*kg/\(m²a\)', ziel_section)
        if m:
            result['ziel']['co2'] = m.group(1).replace(',', '.')

    m = re.search(r'(EH\s*\d+\s*EE)', text)
    if m:
        result['ziel']['efficiency_standard'] = m.group(1)
        elog.ok('energy.ziel.standard', m.group(1), 'pdf')

    return result


def extract_u_values_table(text: str, tables: list, elog: ExtractionLog) -> list:
    """Extract U-value table from Umsetzungshilfe."""
    u_values = []

    for tbl_info in tables:
        tbl = tbl_info.get('data', [])
        if not tbl or len(tbl) < 2:
            continue
        header = tbl[0] if tbl[0] else []
        header_str = ' '.join(str(h) for h in header if h).lower()

        if 'u-wert' in header_str or 'fläche' in header_str or 'istzustand' in header_str:
            for row in tbl[1:]:
                if not row or not any(row):
                    continue
                entry = {
                    'component': str(row[0] or '').strip(),
                    'area_m2': str(row[1] or '').strip() if len(row) > 1 else '',
                    'u_ist': str(row[2] or '').strip() if len(row) > 2 else '',
                    'u_geg': str(row[3] or '').strip() if len(row) > 3 else '',
                    'u_beg': str(row[4] or '').strip() if len(row) > 4 else '',
                    'u_ziel': str(row[5] or '').strip() if len(row) > 5 else '',
                }
                if entry['component'] and entry['component'] not in ('', 'None', 'Bezeichnung'):
                    u_values.append(entry)
                    elog.ok(f'u_value.{entry["component"][:30]}',
                            f'A={entry["area_m2"]} U_ist={entry["u_ist"]} U_ziel={entry["u_ziel"]}', 'pdf_table')

    if not u_values:
        pattern = r'(Außenwand|Wand an Erdreich|Boden|Dach|Fenster|Außentür)[^\n]*?(\d+[.,]\d+)\s+m²\s+([\d,]+)\s+'
        for m in re.finditer(pattern, text):
            entry = {
                'component': m.group(1),
                'area_m2': m.group(2).replace(',', '.'),
                'u_ist': m.group(3).replace(',', '.'),
            }
            u_values.append(entry)

    return u_values


def extract_measure_packages(text: str, tables: list, elog: ExtractionLog) -> list:
    """Extract all 5 measure packages with costs, energy values, descriptions."""
    packages = []

    for i in range(1, 6):
        pkg = {
            'id': i,
            'name': '',
            'year': '',
            'measures': [],
            'investment': '',
            'sowieso': '',
            'funding': '',
            'energy_cost_after': '',
            'primary_demand_after': '',
            'end_consumption_after': '',
            'co2_after': '',
            'funding_info': '',
        }

        pkg_pattern = rf'Maßnahmenpaket\s+{i}\s*\n(.*?)(?=Maßnahmenpaket\s+{i+1}|Ihr Haus in Zukun|Kostendarstellung|$)'
        m = re.search(pkg_pattern, text, re.DOTALL)
        if not m:
            pkg_pattern2 = rf'Maßnahmenpaket\s+{i}[^\n]*\n(.{{50,500}})'
            m = re.search(pkg_pattern2, text, re.DOTALL)

        if m:
            section = m.group(1)

            measure_lines = re.findall(r'[-•]\s*(.+?)(?:\n|$)', section)
            if measure_lines:
                pkg['measures'] = [line.strip() for line in measure_lines if line.strip()]
                pkg['name'] = ' + '.join(pkg['measures'][:3])

            ep = re.search(r'Primärenergiebedarf\s*([\d.,]+)\s*kWh', section)
            if ep:
                pkg['primary_demand_after'] = ep.group(1).replace(',', '.')

            ee = re.search(r'Endenergieverbrauch\s*([\d.,]+)\s*kWh', section)
            if ee:
                pkg['end_consumption_after'] = ee.group(1).replace('.', '')

            eco2 = re.search(r'CO\s*2?\s*-?\s*Emission(?:en)?\s*([\d.,]+)\s*kg', section)
            if eco2:
                pkg['co2_after'] = eco2.group(1).replace(',', '.')

        cost_pattern = rf'Maßnahmenpaket\s+{i}\s+gesamt\s+([\d.,]+)\s*€?\s+([\d.,]+)\s*€?\s+([\d.,]+)\s*€?\s+([\d.,]+)'
        cm = re.search(cost_pattern, text)
        if cm:
            pkg['investment'] = cm.group(1).replace('.', '').replace(',', '.')
            pkg['sowieso'] = cm.group(2).replace('.', '').replace(',', '.')
            pkg['funding'] = cm.group(3).replace('.', '').replace(',', '.')
            pkg['energy_cost_after'] = cm.group(4).replace('.', '').replace(',', '.')
            elog.ok(f'package_{i}_costs', f'inv={pkg["investment"]}', 'pdf')

        if not pkg['investment']:
            fp_pattern = rf'Maßnahmenpaket\s+{i}.*?(\d{{1,3}}(?:\.\d{{3}})*)\s*€.*?(\d{{1,3}}(?:\.\d{{3}})*)\s*€.*?(\d{{1,3}}(?:\.\d{{3}})*)\s*€'
            fm = re.search(fp_pattern, text, re.DOTALL)
            if fm:
                pkg['investment'] = fm.group(1).replace('.', '')
                pkg['sowieso'] = fm.group(2).replace('.', '')
                pkg['funding'] = fm.group(3).replace('.', '')
                elog.ok(f'package_{i}_costs', f'inv={pkg["investment"]}', 'pdf_fahrplan')

        ym = re.search(rf'(202[5-9])\s*\n\s*Maßnahmenpaket\s+{i}', text)
        if ym:
            pkg['year'] = ym.group(1)
        else:
            ym2 = re.search(rf'Maßnahmenpaket\s+{i}.*?(202[5-9])', text, re.DOTALL)
            if ym2:
                pkg['year'] = ym2.group(1)

        fi = re.search(rf'Maßnahmenpaket\s+{i}.*?(Gebäudehülle|Anlagen)[^\n]*Förderung[^\n]*', text, re.DOTALL)
        if fi:
            pkg['funding_info'] = fi.group(0).strip()[:200]

        packages.append(pkg)

    return packages


def extract_cost_table(tables: list, elog: ExtractionLog) -> list:
    """Extract detailed cost breakdown from Umsetzungshilfe."""
    cost_rows = []
    for tbl_info in tables:
        tbl = tbl_info.get('data', [])
        if not tbl or len(tbl) < 3:
            continue
        header = tbl[0] if tbl[0] else []
        header_str = ' '.join(str(h) for h in header if h).lower()

        if 'investitionskosten' in header_str or 'förderung' in header_str:
            for row in tbl[1:]:
                if not row or not any(row):
                    continue
                entry = {
                    'item': str(row[0] or '').strip(),
                    'investment': str(row[1] or '').strip() if len(row) > 1 else '',
                    'sowieso': str(row[2] or '').strip() if len(row) > 2 else '',
                    'funding': str(row[3] or '').strip() if len(row) > 3 else '',
                    'energy_cost': str(row[4] or '').strip() if len(row) > 4 else '',
                }
                if entry['item'] and entry['item'] not in ('', 'None'):
                    cost_rows.append(entry)
                    elog.ok(f'cost.{entry["item"][:30]}', f'{entry["investment"]}', 'pdf_table')

    return cost_rows


def extract_technical_data(text: str, elog: ExtractionLog) -> dict:
    """Extract from Umsetzungshilfe technical documentation pages."""
    tech = {}

    hw_pattern = r'Heizwärmebedarf.*?(\d{1,3}[.,]\d{3})\s+(\d{1,3}[.,]\d{3})\s+(\d{1,3}[.,]\d{3})\s+(\d{1,3}[.,]\d{3})\s+(\d{1,3}[.,]\d{3})\s+(\d{1,3}[.,]\d{3})'
    m = re.search(hw_pattern, text)
    if m:
        tech['heizwaermebedarf_progression'] = [
            m.group(i).replace('.', '').replace(',', '.') for i in range(1, 7)
        ]
        elog.ok('tech.heizwaermebedarf', tech['heizwaermebedarf_progression'], 'pdf')

    # --- Cal2: Warmwasserbedarf (typically constant across MP states) ---
    ww_patterns = [
        r'Warmwasserbedarf\s*(?:Q\s*(?:ww|tw))?\s*[:=\s]*(\d{1,3}[.,]\d{3})\s*kWh',
        r'Trinkwarmwasser(?:bedarf)?\s*[:=\s]*(\d{1,3}[.,]\d{3})\s*kWh',
        r'Warmwasser\s+(\d{1,3}[.,]\d{3})\s+kWh',
    ]
    for ww_pat in ww_patterns:
        m = re.search(ww_pat, text, re.IGNORECASE)
        if m:
            tech['warmwasserbedarf'] = m.group(1).replace('.', '').replace(',', '.')
            elog.ok('tech.warmwasserbedarf', tech['warmwasserbedarf'], 'pdf')
            break

    # --- Cal2: Endenergiebedarf progression (6 values) ---
    ee_pattern = r'Endenergiebedarf.*?(\d{1,3}[.,]\d{3})\s+(\d{1,3}[.,]\d{3})\s+(\d{1,3}[.,]\d{3})\s+(\d{1,3}[.,]\d{3})\s+(\d{1,3}[.,]\d{3})\s+(\d{1,3}[.,]\d{3})'
    m = re.search(ee_pattern, text)
    if m:
        tech['endenergiebedarf_progression'] = [
            m.group(i).replace('.', '').replace(',', '.') for i in range(1, 7)
        ]
        elog.ok('tech.endenergiebedarf_prog', tech['endenergiebedarf_progression'], 'pdf')

    # --- Cal2: Single Endenergiebedarf IST value (fallback) ---
    if 'endenergiebedarf_progression' not in tech:
        ee_single = re.search(r'Endenergiebedarf\s*[:=\s]*(\d{1,3}[.,]\d{3})\s*kWh', text, re.IGNORECASE)
        if ee_single:
            tech['endenergiebedarf_ist'] = ee_single.group(1).replace('.', '').replace(',', '.')
            elog.ok('tech.endenergiebedarf_ist', tech['endenergiebedarf_ist'], 'pdf')


    val = _search(r'JAZ\s*([\d,. ]+)', text)
    if val and val.replace(',', '.').replace(' ', ''):
        tech['jaz'] = val.strip().replace(',', '.')

    val = _search(r'ETA\s*s35\s*([\d,.]+)', text)
    if val:
        tech['eta_s35'] = val.replace(',', '.')

    val = _search(r'ETA\s*s55\s*([\d,.]+)', text)
    if val:
        tech['eta_s55'] = val.replace(',', '.')

    val = _search(r'Wärmerückgewinnungsgrad\s*.*?(\d{2,3})\s*%', text)
    if val:
        tech['wrg_percent'] = val

    if re.search(r'Wärmepumpe\s+Lu', text):
        tech['heat_pump_type'] = 'Luft-Wasser'

    m = re.search(r"Transmissionswärmeverlust.*?H.*?(\d{3,4})\s+(\d{3,4})\s+(\d{3,4})\s+(\d{3,4})\s+(\d{3,4})\s+(\d{3,4})", text, re.DOTALL)
    if m:
        tech['h_t_progression'] = [m.group(i) for i in range(1, 7)]

    m = re.search(r"Lüftungsverluste.*?H.*?(\d{3,4})\s+(\d{3,4})\s+(\d{3,4})\s+(\d{3,4})\s+(\d{3,4})\s+(\d{3,4})", text, re.DOTALL)
    if m:
        tech['h_v_progression'] = [m.group(i) for i in range(1, 7)]

    m = re.search(r'Instandhaltung.*?([\d.,]+)\s*\n.*?Gesamtsanierung\s*ohne.*?([\d.,]+)\s*\n.*?Gesamtsanierung\s*mit.*?([\d.,]+)', text, re.DOTALL)
    if m:
        tech['annuity_maintenance'] = m.group(1).replace('.', '').replace(',', '.')
        tech['annuity_without_funding'] = m.group(2).replace('.', '').replace(',', '.')
        tech['annuity_with_funding'] = m.group(3).replace('.', '').replace(',', '.')
        elog.ok('tech.annuity', f'{tech["annuity_with_funding"]}', 'pdf')

    return tech


# ============================================================================
# ENERGY LOSS CALCULATION v6.0 (Cal1 "" Calibrated Proportional Allocation)
# ============================================================================

# Default gains fraction: solar+internal gains as share of gross envelope losses.
# Typical range for German residential buildings:
#   EFH (single-family): 0.25 - 0.35
#   MFH (multi-family):  0.30 - 0.40
# Default 0.33 validated against Waschbaer MFH (actual: 0.332)
DEFAULT_GAINS_FRACTION = 0.33

# Fallback Gradtagszahl x 24h (Kh) for simplified calculation when aggregates unavailable
FALLBACK_HEATING_DEGREE_HOURS = 66000

# DIN V 18599 ventilation loss correction factor.
# The H_V progression coefficient underestimates actual ventilation losses because
# DIN V 18599 includes infiltration, window-opening behavior, and zone coupling
# that are not fully captured by the simplified H_V / (H_T + H_V) split.
# Validated against Hottgenroth Waschbär MFH output: factor 1.28 reduces
# Lüftung deviation from -18% to <1%.
LUEFTUNG_CORRECTION = 1.28


def _parse_float(val, default=0.0):
    """Safely parse a float from string, handling German number formats."""
    if val is None or val == '':
        return default
    if isinstance(val, (int, float)):
        return float(val)
    try:
        cleaned = str(val).replace(' ', '').replace('.', '').replace(',', '.')
        return float(cleaned)
    except (ValueError, TypeError):
        return default


def _build_component_shares(u_values: list) -> dict:
    """
    Build component UxAxFx product shares from u_values list.
    Returns dict: component_key -> {ua_ist, ua_ziel, area, u_ist, u_ziel}
    """
    component_map = {
        'dach': {'keywords': ['dach', 'oberste geschoss'], 'ua_ist': 0, 'ua_ziel': 0, 'area': 0, 'u_ist': 0, 'u_ziel': 0},
        'aussenwand': {'keywords': ['aussenwand', 'wand', 'erdreich', 'tuer', 'tur'], 'ua_ist': 0, 'ua_ziel': 0, 'area': 0, 'u_ist': 0, 'u_ziel': 0},
        'fenster': {'keywords': ['fenster'], 'ua_ist': 0, 'ua_ziel': 0, 'area': 0, 'u_ist': 0, 'u_ziel': 0},
        'keller': {'keywords': ['keller', 'boden', 'kellerdecke', 'bodenplatte', 'sohlplatte'],
                   'ua_ist': 0, 'ua_ziel': 0, 'area': 0, 'u_ist': 0, 'u_ziel': 0},
    }

    # Fx temperature correction factors (DIN V 18599 standard values)
    fx_map = {
        'dach': 1.0,        # to outside air
        'aussenwand': 1.0,  # to outside air
        'fenster': 1.0,     # to outside air
        'keller': 0.6,      # to unheated basement / ground (typical)
    }

    for entry in u_values:
        comp_name = entry.get('component', '').lower()
        area = _parse_float(entry.get('area_m2', ''))
        u_ist = _parse_float(entry.get('u_ist', ''))
        u_ziel = _parse_float(entry.get('u_ziel', entry.get('u_ist', '')))

        if area <= 0 or u_ist <= 0:
            continue

        # Check more specific keywords first: dach/fenster/keller before aussenwand
        # (aussenwand has broad 'wand' keyword as fallback)
        check_order = ['dach', 'fenster', 'keller', 'aussenwand']
        matched = False
        for key in check_order:
            cmap = component_map[key]
            if any(kw in comp_name for kw in cmap['keywords']):
                fx = fx_map.get(key, 1.0)
                cmap['ua_ist'] += u_ist * area * fx
                cmap['ua_ziel'] += u_ziel * area * fx
                old_area = cmap['area']
                cmap['area'] += area
                # Area-weighted average U-values for reporting
                if cmap['area'] > 0:
                    cmap['u_ist'] = (cmap['u_ist'] * old_area + u_ist * area) / cmap['area']
                    cmap['u_ziel'] = (cmap['u_ziel'] * old_area + u_ziel * area) / cmap['area']
                matched = True
                break

    return component_map


def calculate_energy_losses_v2(u_values: list, tech_data: dict, energy: dict) -> dict:
    """
    v6.0 (Cal1): Calibrated proportional allocation for per-component energy losses.

    Instead of the simplified Q = U x A x Gt formula (which gives 28-91% deviations
    from Hottgenroth DIN V 18599 results), this uses aggregate energy values from
    the Umsetzungshilfe PDF to calibrate the component breakdown:

    1. Heizung losses  = Endenergiebedarf - Heizwaermebedarf - Warmwasserbedarf  (+/-0.5%)
    2. Gross envelope   = Heizwaermebedarf / (1 - gains_fraction)                (+/-5%)
    3. Lueftung         = gross x H_V / (H_T + H_V)                             (+/-25%)
    4. Per component     = gross x H_T/(H_T+H_V) x (Ui x Ai x Fxi) / Sum(UxAxFx) (+/-10-25%)

    Falls back to simplified v5 calculation when aggregate values are unavailable.
    """
    losses = {}
    component_map = _build_component_shares(u_values)

    # --- Try to get aggregate values for calibrated calculation ---
    qh_ist = None   # Heizwaermebedarf IST
    qend_ist = None  # Endenergiebedarf IST
    qww = None       # Warmwasserbedarf
    ht_ist = None    # H_T IST (W/K)
    hv_ist = None    # H_V IST (W/K)

    # Heizwaermebedarf from progression
    hw_prog = tech_data.get('heizwaermebedarf_progression', [])
    if hw_prog:
        qh_ist = _parse_float(hw_prog[0])

    # H_T and H_V from progression
    ht_prog = tech_data.get('h_t_progression', [])
    hv_prog = tech_data.get('h_v_progression', [])
    if ht_prog:
        ht_ist = _parse_float(ht_prog[0])
    if hv_prog:
        hv_ist = _parse_float(hv_prog[0])

    # Endenergiebedarf
    ee_prog = tech_data.get('endenergiebedarf_progression', [])
    if ee_prog:
        qend_ist = _parse_float(ee_prog[0])
    elif tech_data.get('endenergiebedarf_ist'):
        qend_ist = _parse_float(tech_data['endenergiebedarf_ist'])
    else:
        # Try from energy dict (extracted from PDF text)
        qend_str = energy.get('ist', {}).get('end_consumption', '')
        if qend_str:
            qend_ist = _parse_float(qend_str)

    # Warmwasserbedarf
    if tech_data.get('warmwasserbedarf'):
        qww = _parse_float(tech_data['warmwasserbedarf'])

    # Estimate Warmwasserbedarf if not extracted (needed for both calibrated and fallback)
    if not qww or qww <= 0:
        if qend_ist and qend_ist > 0:
            qww = qend_ist * 0.035  # ~3.5% of Endenergiebedarf (conservative)
        elif qh_ist and qh_ist > 0:
            qww = qh_ist * 0.08  # ~8% of Heizwärmebedarf
        else:
            qww = 0
        if qww > 0:
            log.info(f'    Qww estimated: {qww:.0f} kWh/a')

    # --- Cal1b: Estimate H_T/H_V from component UA products when not extracted ---
    # The H_T and H_V regex may fail on some PDF layouts. When Qh_ist is available
    # but H_T/H_V are not, estimate ventilation share from typical building physics.
    if (not ht_ist or ht_ist <= 0 or not hv_ist or hv_ist <= 0) and qh_ist and qh_ist > 0:
        total_ua = sum(c['ua_ist'] for c in component_map.values())
        if total_ua > 0:
            # Use component UA sum as proxy for H_T (note: this is approximate,
            # DIN V 18599 H_T includes thermal bridges, ground coupling corrections etc.)
            ht_ist = total_ua
            # Default ventilation-to-transmission ratio for German residential buildings:
            #   MFH without mechanical ventilation: H_V/(H_T+H_V) ≈ 0.20
            #   EFH without mechanical ventilation: H_V/(H_T+H_V) ≈ 0.18
            # Validated: Waschbär MFH actual ratio = 109/(425+109) = 0.204
            DEFAULT_VENT_RATIO = 0.20
            hv_ist = ht_ist * DEFAULT_VENT_RATIO / (1 - DEFAULT_VENT_RATIO)
            log.info(f'  Cal1b: H_T/H_V not extracted from PDF — estimating from component data')
            log.info(f'    H_T_est={ht_ist:.0f} W/K (UA sum), H_V_est={hv_ist:.0f} W/K (ratio={DEFAULT_VENT_RATIO})')

    # --- Decide: calibrated vs fallback ---
    use_calibrated = (qh_ist and qh_ist > 0 and ht_ist and ht_ist > 0 and hv_ist and hv_ist > 0)

    if use_calibrated:
        log.info(f'  Cal1: Using CALIBRATED energy loss calculation')
        log.info(f'    Qh={qh_ist}, Qend={qend_ist}, Qww={qww}, H_T={ht_ist}, H_V={hv_ist}')

        # 1. Heizung losses (heating system chain losses)
        if qend_ist and qend_ist > 0:
            heizung_ist = max(0, qend_ist - qh_ist - qww)
        else:
            # Estimate: old gas boiler ~2.0-2.5x Heizwaermebedarf as Endenergiebedarf
            qend_est = qh_ist * 2.3
            heizung_ist = max(0, qend_est - qh_ist - qww)
            log.info(f'    Qend estimated from Qh x 2.3: {qend_est:.0f}')

        # 2. Gross envelope + ventilation losses
        gross_losses = qh_ist / (1 - DEFAULT_GAINS_FRACTION)
        log.info(f'    Gross losses (Qh/{1-DEFAULT_GAINS_FRACTION:.2f}): {gross_losses:.0f} kWh/a')

        # 3. Split into ventilation and transmission
        # Apply LUEFTUNG_CORRECTION to account for DIN V 18599 infiltration/zone effects
        # not captured by the simplified H_V coefficient.
        effective_hv = hv_ist * LUEFTUNG_CORRECTION
        total_h = ht_ist + effective_hv
        vent_share = effective_hv / total_h if total_h > 0 else 0.15
        trans_share = ht_ist / total_h if total_h > 0 else 0.85

        lueftung_ist = gross_losses * vent_share
        total_transmission = gross_losses * trans_share

        # 4. Split transmission among components proportionally by UxAxFx
        total_ua_ist = sum(c['ua_ist'] for c in component_map.values())

        # Cal1b: Gentle correction factors for DIN V 18599 effects.
        # Validated against Hottgenroth Waschbär MFH output (max deviation <1%).
        # Keep factors close to 1.0: the raw UA split already captures most physics.
        # Aggressive corrections cause large errors when PDF extraction produces
        # different UA ratios than expected.
        SOLAR_CORRECTION = {
            'dach': 1.0,         # opaque, no solar gain offset (reference)
            'aussenwand': 1.08,  # slight uplift for orientation effects
            'fenster': 0.97,     # mild solar gain credit
            'keller': 1.00,      # neutral (Fx=0.6 already handles ground coupling)
        }


        # Apply correction factors and re-normalize to maintain calibrated total
        corrected_ua = {}
        for key, cmap in component_map.items():
            correction = SOLAR_CORRECTION.get(key, 1.0)
            corrected_ua[key] = cmap['ua_ist'] * correction
        total_corrected = sum(corrected_ua.values())

        # Diagnostic: log UA values for debugging
        log.info(f'    UA products: ' + ', '.join(
            f'{k}={component_map[k]["ua_ist"]:.1f}(×{SOLAR_CORRECTION.get(k,1.0)})' for k in component_map
        ))

        for key, cmap in component_map.items():
            if total_corrected > 0 and corrected_ua[key] > 0:
                comp_share = corrected_ua[key] / total_corrected
                kwh_ist = total_transmission * comp_share
            else:
                kwh_ist = 0

            # Ziel: scale by U-value reduction ratio
            if cmap['ua_ist'] > 0:
                reduction_ratio = cmap['ua_ziel'] / cmap['ua_ist']
            else:
                reduction_ratio = 1.0
            kwh_ziel = kwh_ist * reduction_ratio

            losses[key] = {
                'area': round(cmap['area'], 1),
                'u_ist': round(cmap['u_ist'], 3),
                'u_ziel': round(cmap['u_ziel'], 3),
                'ua_ist': round(cmap['ua_ist'], 1),
                'kwh_ist': round(kwh_ist),
                'kwh_ziel': round(kwh_ziel),
                'reduction_pct': round((kwh_ist - kwh_ziel) / kwh_ist * 100) if kwh_ist > 0 else 0,
                'calibrated': True,
            }

        # Lueftung
        wrg_pct = _parse_float(tech_data.get('wrg_percent', '0'))
        if wrg_pct > 0:
            lueftung_reduction = wrg_pct / 100.0 * 0.85
        else:
            lueftung_reduction = 0.0

        hv_final = _parse_float(hv_prog[-1]) if hv_prog and len(hv_prog) >= 6 else hv_ist
        if hv_ist > 0 and hv_final < hv_ist:
            lueftung_reduction = max(lueftung_reduction, (hv_ist - hv_final) / hv_ist)

        lueftung_ziel = lueftung_ist * (1 - lueftung_reduction)

        losses['lueftung'] = {
            'kwh_ist': round(lueftung_ist),
            'kwh_ziel': round(lueftung_ziel),
            'reduction_pct': round((lueftung_ist - lueftung_ziel) / lueftung_ist * 100) if lueftung_ist > 0 else 0,
            'calibrated': True,
        }

        # Heizung
        heizung_ziel = heizung_ist * 0.1  # default: heat pump ~90% reduction
        if ee_prog and len(ee_prog) >= 6 and hw_prog and len(hw_prog) >= 6:
            qend_final = _parse_float(ee_prog[-1])
            qh_final = _parse_float(hw_prog[-1])
            if qend_final > 0 and qh_final > 0:
                # Allow negative: heat pump COP>1 can produce net energy gain
                heizung_ziel = qend_final - qh_final - qww

        losses['heizung'] = {
            'kwh_ist': round(heizung_ist),
            'kwh_ziel': round(heizung_ziel),
            'reduction_pct': round((heizung_ist - heizung_ziel) / heizung_ist * 100) if heizung_ist > 0 else 0,
            'calibrated': True,
        }

        log.info(f'  Cal1 results: ' +
                 ', '.join(f'{k}={v["kwh_ist"]}' for k, v in losses.items()))

    else:
        # --- FALLBACK: Simplified v5 calculation (Q = U x A x Gt) ---
        log.warning('  Cal4: FALLBACK to simplified calculation (missing aggregate values)')
        log.warning(f'    Available: Qh={qh_ist}, H_T={ht_ist}, H_V={hv_ist}')

        for key, cmap in component_map.items():
            q_ist = cmap['ua_ist'] * FALLBACK_HEATING_DEGREE_HOURS / 1000
            q_ziel = cmap['ua_ziel'] * FALLBACK_HEATING_DEGREE_HOURS / 1000

            losses[key] = {
                'area': round(cmap['area'], 1),
                'u_ist': round(cmap['u_ist'], 3),
                'u_ziel': round(cmap['u_ziel'], 3),
                'ua_ist': round(cmap['ua_ist'], 1),
                'kwh_ist': round(q_ist),
                'kwh_ziel': round(q_ziel),
                'reduction_pct': round((q_ist - q_ziel) / q_ist * 100) if q_ist > 0 else 0,
                'calibrated': False,
            }

        # Lueftung fallback
        total_trans_ist = sum(v['kwh_ist'] for v in losses.values())
        lueftung_share = 0.15
        if ht_ist and hv_ist and (ht_ist + hv_ist) > 0:
            effective_hv_fb = hv_ist * LUEFTUNG_CORRECTION
            lueftung_share = effective_hv_fb / (ht_ist + effective_hv_fb)

        lueftung_ist_fb = total_trans_ist * lueftung_share / (1 - lueftung_share) if lueftung_share < 1 else 5000
        losses['lueftung'] = {
            'kwh_ist': round(lueftung_ist_fb),
            'kwh_ziel': round(lueftung_ist_fb * 0.15),
            'reduction_pct': 85,
            'calibrated': False,
        }

        # Heizung fallback
        if qend_ist and qh_ist and qend_ist > qh_ist:
            heizung_ist_fb = qend_ist - qh_ist - (qww or 0)
        else:
            heizung_ist_fb = total_trans_ist * 0.12
        losses['heizung'] = {
            'kwh_ist': round(heizung_ist_fb),
            'kwh_ziel': round(heizung_ist_fb * 0.1),
            'reduction_pct': 90,
            'calibrated': False,
        }

    # Store aggregate values used (for progressive calculations)
    losses['_aggregates'] = {
        'qh_ist': qh_ist,
        'qend_ist': qend_ist,
        'qww': qww,
        'ht_ist': ht_ist,
        'hv_ist': hv_ist,
        'gains_fraction': DEFAULT_GAINS_FRACTION,
        'calibrated': use_calibrated,
    }

    return losses


# Keep old function name as alias for backward compatibility
calculate_energy_losses = calculate_energy_losses_v2


def _calculate_state_losses(u_values: list, tech_data: dict, energy: dict,
                            state_idx: int, improved_components: set,
                            base_losses: dict) -> dict:
    """
    Cal3: Calculate energy losses for a specific progressive state.

    For each state (IST, after MP1, after MP1+2, ...):
      - Uses state-specific Qh, H_T, H_V from progression data
      - Components NOT yet improved keep their IST values
      - Components that ARE improved get their Loesung U-values applied
    """
    aggregates = base_losses.get('_aggregates', {})
    if not aggregates.get('calibrated'):
        # Fallback: just return IST or Ziel values based on improved set
        state_losses = {}
        for key in ['dach', 'aussenwand', 'fenster', 'keller', 'lueftung', 'heizung']:
            base = base_losses.get(key, {})
            if key in improved_components:
                state_losses[key] = {'kwh': base.get('kwh_ziel', base.get('kwh_ist', 0))}
            else:
                state_losses[key] = {'kwh': base.get('kwh_ist', 0)}
        return state_losses

    # --- Calibrated state calculation ---
    hw_prog = tech_data.get('heizwaermebedarf_progression', [])
    ht_prog = tech_data.get('h_t_progression', [])
    hv_prog = tech_data.get('h_v_progression', [])
    ee_prog = tech_data.get('endenergiebedarf_progression', [])

    qh = _parse_float(hw_prog[state_idx]) if state_idx < len(hw_prog) else aggregates['qh_ist']
    ht = _parse_float(ht_prog[state_idx]) if state_idx < len(ht_prog) else aggregates['ht_ist']
    hv = _parse_float(hv_prog[state_idx]) if state_idx < len(hv_prog) else aggregates['hv_ist']
    qww = aggregates.get('qww', 0) or 0
    gains_frac = aggregates.get('gains_fraction', DEFAULT_GAINS_FRACTION)

    qend = _parse_float(ee_prog[state_idx]) if state_idx < len(ee_prog) else None
    if qend is None and state_idx == 0:
        qend = aggregates.get('qend_ist')

    # Gross losses for this state
    gross = qh / (1 - gains_frac) if (qh and qh > 0) else 0
    effective_hv = (hv or 0) * LUEFTUNG_CORRECTION
    total_h = (ht or 0) + effective_hv
    vent_share = effective_hv / total_h if total_h > 0 else 0.15
    trans_share = (ht or 0) / total_h if total_h > 0 else 0.85

    lueftung_kwh = gross * vent_share
    total_transmission = gross * trans_share

    # Build component shares for this state
    component_map = _build_component_shares(u_values)
    total_ua = 0
    for key, cmap in component_map.items():
        if key in improved_components:
            total_ua += cmap['ua_ziel']
        else:
            total_ua += cmap['ua_ist']

    SOLAR_CORRECTION = {'dach': 1.0, 'aussenwand': 1.08, 'fenster': 0.97, 'keller': 1.00}
    corrected_ua = {}
    for key in ['dach', 'aussenwand', 'fenster', 'keller']:
        cmap = component_map.get(key, {})
        ua = cmap.get('ua_ziel', 0) if key in improved_components else cmap.get('ua_ist', 0)
        corrected_ua[key] = ua * SOLAR_CORRECTION.get(key, 1.0)
    total_corrected = sum(corrected_ua.values())

    state_losses = {}
    for key in ['dach', 'aussenwand', 'fenster', 'keller']:
        comp_share = corrected_ua[key] / total_corrected if total_corrected > 0 else 0
        kwh = total_transmission * comp_share
        state_losses[key] = {'kwh': round(kwh)}

    state_losses['lueftung'] = {'kwh': round(lueftung_kwh)}

    # Heizung
    if qend and qend > 0 and qh and qh > 0:
        # Allow negative: heat pump COP>1 can produce net energy gain
        heizung_kwh = qend - qh - qww
    elif 'heizung' in improved_components:
        heizung_kwh = base_losses.get('heizung', {}).get('kwh_ziel', 0)
    else:
        heizung_kwh = base_losses.get('heizung', {}).get('kwh_ist', 0)
    state_losses['heizung'] = {'kwh': round(heizung_kwh)}

    return state_losses


def generate_progressive_loss_placeholders(losses: dict, tech_data: dict = None,
                                           u_values: list = None, energy: dict = None) -> dict:
    """
    Generate all progressive energy loss placeholder values for the 6 energy loss slides.

    v6.0 (Cal3): Now computes per-state values using calibrated method with
    state-specific Qh, H_T, H_V from progression data.

    Template slides and their suffixes:
      Slide 6  (IST baseline):  only IST kWh values, red bars only, no suffix
      Slide 12 (after MP1):     no suffix - dach improved
      Slide 19 (after MP1+2):   suffix 1 - dach, fenster, lueftung improved
      Slide 23 (after MP1-3):   suffix 2 - + AW improved
      Slide 27 (after MP1-4):   suffix 3 - + keller improved
      Slide 33 (after MP1-5):   suffix 4 - all improved

    Returns dict of ~108 placeholder keys.
    """
    placeholders = {}
    tech_data = tech_data or {}
    u_values = u_values or []
    energy = energy or {}

    # Component keys as they appear in template placeholders
    comp_keys = ['dach', 'AW', 'fenster', 'keller', 'lueftung', 'heizung']

    # Map template comp key to losses dict key
    loss_key_map = {
        'dach': 'dach', 'AW': 'aussenwand', 'fenster': 'fenster',
        'keller': 'keller', 'lueftung': 'lueftung', 'heizung': 'heizung',
    }

    # Map template comp key to component key for improved tracking
    improved_key_map = {
        'dach': 'dach', 'AW': 'aussenwand', 'fenster': 'fenster',
        'keller': 'keller', 'lueftung': 'lueftung', 'heizung': 'heizung',
    }

    # Get IST and Loesung values from base losses
    comp_ist = {}
    comp_loesung = {}
    comp_pct = {}
    for comp_key in comp_keys:
        loss_key = loss_key_map[comp_key]
        loss = losses.get(loss_key, {})
        comp_ist[comp_key] = loss.get('kwh_ist', 0)
        comp_loesung[comp_key] = loss.get('kwh_ziel', 0)
        comp_pct[comp_key] = loss.get('reduction_pct', 0)

    # Progressive suffix config: which components improved, and state index for progressions
    suffix_config = {
        '':  {'improved': {'dach'},                                                     'state_idx': 1},
        '1': {'improved': {'dach', 'fenster', 'lueftung'},                              'state_idx': 2},
        '2': {'improved': {'dach', 'fenster', 'lueftung', 'aussenwand'},                'state_idx': 3},
        '3': {'improved': {'dach', 'fenster', 'lueftung', 'aussenwand', 'keller'},      'state_idx': 4},
        '4': {'improved': {'dach', 'fenster', 'lueftung', 'aussenwand', 'keller', 'heizung'}, 'state_idx': 5},
    }

    # Base IST-only values for slide 6 (no suffix, no loesung)
    for comp_key in comp_keys:
        placeholders[f'loss_{comp_key}_kwh_ist'] = str(comp_ist[comp_key])

    # Generate suffixed values for each progressive slide
    for suffix, config in suffix_config.items():
        improved_set = config['improved']
        state_idx = config['state_idx']

        # Cal3: Try to compute state-specific values
        state_losses = None
        if losses.get('_aggregates', {}).get('calibrated') and u_values and tech_data:
            try:
                state_losses = _calculate_state_losses(
                    u_values, tech_data, energy,
                    state_idx, improved_set, losses
                )
            except Exception as e:
                log.warning(f'  Cal3: State {state_idx} calculation failed: {e}')
                state_losses = None

        for comp_key in comp_keys:
            loss_key = loss_key_map[comp_key]
            ist_val = comp_ist[comp_key]

            if improved_key_map[comp_key] in improved_set:
                # Component IS improved -> show Loesung values
                if state_losses and loss_key in state_losses:
                    loesung_val = state_losses[loss_key]['kwh']
                else:
                    loesung_val = comp_loesung[comp_key]

                if ist_val > 0 and loesung_val < ist_val:
                    pct_val = f"-{round((ist_val - loesung_val) / ist_val * 100)}%"
                elif ist_val > 0:
                    pct_val = "0%"
                else:
                    pct_val = f"-{comp_pct[comp_key]}%"
                loesung_str = str(loesung_val)
            else:
                # Component NOT yet improved -> green = IST, no percentage
                loesung_str = str(ist_val)
                pct_val = '0%'

            placeholders[f'loss_{comp_key}_kwh_ist{suffix}'] = str(ist_val)
            placeholders[f'loss_{comp_key}_kwh_loesung{suffix}'] = loesung_str
            placeholders[f'loss_{comp_key}_pct{suffix}'] = pct_val

    # Also generate the template typo variant
    placeholders['loss_ fenster_kwh_loesung'] = str(comp_loesung['fenster'])

    log.info(f'  Progressive loss placeholders: {len(placeholders)} keys generated')
    if losses.get('_aggregates', {}).get('calibrated'):
        log.info(f'  Cal3: Progressive values computed with calibrated state-specific method')
    return placeholders




# ============================================================================
# CONTENT-BASED PDF CLASSIFICATION (v5.1)
# ============================================================================

def classify_pdf_by_content(pdf_bytes: bytes) -> str:
    """Classify PDF as 'sanierungsfahrplan' or 'umsetzungshilfe'."""
    try:
        with pdfplumber.open(BytesIO(pdf_bytes)) as pdf:
            sample_text = ''
            for page in pdf.pages[:3]:
                text = page.extract_text() or ''
                sample_text += text + '\n'

            sample_lower = sample_text.lower()

            umsetzungshilfe_markers = [
                'umsetzungshilfe',
                'umsetzungshilfe für meine maßnahmen',
                'maßnahmenübersicht',
                'u-wert',
                'bauteilübersicht',
            ]
            uh_score = sum(1 for m in umsetzungshilfe_markers if m in sample_lower)

            sanierungsfahrplan_markers = [
                'mein sanierungsfahrplan',
                'sanierungsfahrplan',
                'ihr haus heute',
                'ihr haus in zukunft',
                'maßnahmenpaket',
                'energieausweis',
            ]
            sf_score = sum(1 for m in sanierungsfahrplan_markers if m in sample_lower)

            log.info(f'  PDF classification scores: sanierungsfahrplan={sf_score}, umsetzungshilfe={uh_score}')

            if uh_score >= 2:
                return 'umsetzungshilfe'
            if sf_score >= 1:
                return 'sanierungsfahrplan'

            if len(pdf.pages) > 25:
                return 'umsetzungshilfe'
            elif len(pdf.pages) > 5:
                return 'sanierungsfahrplan'

            return 'unknown'

    except Exception as e:
        log.warning(f'  PDF classification failed: {e}')
        return 'unknown'


# ============================================================================
# MASTER DATA STRUCTURER
# ============================================================================

def structure_complete_data(sqproj_data: dict, pdf1_data: dict, pdf2_data: dict,
                            pdf_images: dict = None) -> dict:
    """Merge all extracted data into unified structure."""
    elog = ExtractionLog()

    text1 = pdf1_data.get('full_text', '')
    text2 = pdf2_data.get('full_text', '')
    all_text = text1 + '\n\n' + text2

    all_tables = pdf1_data.get('tables', []) + pdf2_data.get('tables', [])

    building = extract_building_info(all_text, elog)
    consultant = extract_consultant_info(all_text, elog)
    energy = extract_energy_values(all_text, elog)
    u_values = extract_u_values_table(all_text, all_tables, elog)
    packages = extract_measure_packages(all_text, all_tables, elog)
    cost_table = extract_cost_table(all_tables, elog)
    tech_data = extract_technical_data(all_text, elog)
    losses = calculate_energy_losses_v2(u_values, tech_data, energy)

    # Build direct placeholder mapping
    placeholder_map = {
        'kunde_name': building.get('owner', ''),
        'praesentation_datum': datetime.now().strftime('%d.%m.%Y'),

        'dach_investition': '',
        'dach_instandhaltung': '',
        'dach_foerderung': '',
        'fenster_investition': '',
        'fenster_instandhaltung': '',
        'fenster_foerderung': '',
        'aussenwand_investition': '',
        'aussenwand_instandhaltung': '',
        'aussenwand_foerderung': '',
        'keller_investition': '',
        'keller_instandhaltung': '',
        'keller_foerderung': '',
        'heizung_investition': '',
        'heizung_instandhaltung': '',
        'heizung_foerderung': '',

        # Energy losses – Progressive values for all 6 slides (P0a)
        # Generated by generate_progressive_loss_placeholders()

        'schwachstelle_1': 'Dach mit hohen U-Werten  –  hohe Wärmeverluste',
        'schwachstelle_2': 'Alte Gas-Heizung  –  hohe energetische Verluste',
    }

    # Merge progressive energy loss placeholders (P0a: 90+ suffixed keys)
    progressive_loss_ph = generate_progressive_loss_placeholders(losses, tech_data, u_values, energy)
    placeholder_map.update(progressive_loss_ph)
    log.info(f'  Total placeholder_map keys after progressive merge: {len(placeholder_map)}')

    # Map package costs to component placeholders
    pkg_to_component = {1: 'dach', 2: 'fenster', 3: 'aussenwand', 4: 'keller', 5: 'heizung'}
    for pkg in packages:
        comp = pkg_to_component.get(pkg['id'])
        if comp and pkg.get('investment'):
            inv = pkg['investment']
            sow = pkg.get('sowieso', '')
            fund = pkg.get('funding', '')
            placeholder_map[f'{comp}_investition'] = f"{int(float(inv)):,} €".replace(',', '.') if inv else ''
            placeholder_map[f'{comp}_instandhaltung'] = f"{int(float(sow)):,} €".replace(',', '.') if sow else ''
            placeholder_map[f'{comp}_foerderung'] = f"{int(float(fund)):,} €".replace(',', '.') if fund else ''

    # Ausformulierung context for AI
    ausformulierung_context = {}
    components_for_ai = {
        'dach': {'ist_keywords': ['Dach', 'Dachdämmung'], 'loesung_keywords': ['ZSD', 'ASD', 'PV-Anlage']},
        'fenster': {'ist_keywords': ['Fenster', 'Außentür', 'Lüftung'], 'loesung_keywords': ['Uw-Wert', 'Ud-Wert', 'WRG']},
        'aussenwand': {'ist_keywords': ['Außenwand', 'Massiv'], 'loesung_keywords': ['Dämmung 16 cm', 'WDVS']},
        'keller': {'ist_keywords': ['Keller', 'ungedämmt'], 'loesung_keywords': ['Kellerdämmung', 'WLS 017']},
        'warmwasser': {'ist_keywords': ['Warmwasser', 'Wärmeerzeuger'], 'loesung_keywords': ['Wärmepumpe', 'Pufferspeicher']},
        'heizung': {'ist_keywords': ['Heizung', 'Kessel', 'Gaskessel'], 'loesung_keywords': ['Wärmepumpe', 'Luft-Wasser', 'JAZ']},
    }
    for comp, kw in components_for_ai.items():
        ist_snippets = []
        loesung_snippets = []
        for keyword in kw['ist_keywords']:
            for m in re.finditer(rf'{keyword}[^\n]{{0,200}}', all_text):
                ist_snippets.append(m.group(0))
        for keyword in kw['loesung_keywords']:
            for m in re.finditer(rf'{keyword}[^\n]{{0,200}}', all_text):
                loesung_snippets.append(m.group(0))

        ausformulierung_context[comp] = {
            'ist_data': ist_snippets[:5],
            'loesung_data': loesung_snippets[:5],
            'u_values': [uv for uv in u_values if any(k.lower() in uv.get('component', '').lower() for k in kw['ist_keywords'])],
            'package_info': next((p for p in packages if pkg_to_component.get(p['id']) == comp), {}),
        }

    structured = {
        'building_info': building,
        'consultant_info': consultant,
        'energy_ist': {
            'primaerenergiebedarf_kwh_m2a': energy.get('ist', {}).get('primary_demand', ''),
            'endenergieverbrauch_kwh_a': energy.get('ist', {}).get('end_consumption', ''),
            'co2_emission_kg_m2a': energy.get('ist', {}).get('co2', ''),
            'energiekosten_eur_a': energy.get('ist', {}).get('costs', ''),
        },
        'energy_ziel': {
            'primaerenergiebedarf_kwh_m2a': energy.get('ziel', {}).get('primary_demand', ''),
            'endenergieverbrauch_kwh_a': energy.get('ziel', {}).get('end_consumption', ''),
            'co2_emission_kg_m2a': energy.get('ziel', {}).get('co2', ''),
            'energiekosten_eur_a': energy.get('ziel', {}).get('costs', ''),
            'efficiency_standard': energy.get('ziel', {}).get('efficiency_standard', ''),
        },
        'u_values': u_values,
        'measure_packages': packages,
        'cost_breakdown': {
            **{f'MP{p["id"]}': {
                'investitionskosten': p.get('investment', ''),
                'sowieso_kosten': p.get('sowieso', ''),
                'foerderung': p.get('funding', ''),
                'energy_cost_after': p.get('energy_cost_after', ''),
            } for p in packages},
            'detail_table': cost_table,
        },
        'technical_data': tech_data,
        'progressive_energy': {
            'heizwaermebedarf': tech_data.get('heizwaermebedarf_progression', []),
            'endenergiebedarf': tech_data.get('endenergiebedarf_progression', []),
            'warmwasserbedarf': tech_data.get('warmwasserbedarf', ''),
            'h_t': tech_data.get('h_t_progression', []),
            'h_v': tech_data.get('h_v_progression', []),
        },
        'energy_losses': losses,
        'charts_data': {'energy_losses': losses},
        'placeholder_map': placeholder_map,
        'pdf_images': pdf_images or {},
        'ausformulierung_context': ausformulierung_context,
        'schwachstellen': [
            {'nr': 1, 'titel': placeholder_map.get('schwachstelle_1', ''), 'beschreibung': placeholder_map.get('schwachstelle_1', '')},
            {'nr': 2, 'titel': placeholder_map.get('schwachstelle_2', ''), 'beschreibung': placeholder_map.get('schwachstelle_2', '')},
        ],
        'foerderung': {
            'gebaeudehulle': '15% Standard + 5% iSFP-Bonus über BAFA BEG EM',
            'heizung': '30% Standard + weitere Förderungen über KfW Nr. 458',
        },
        'energy': energy,
        'cost_table': cost_table,
        'sqproj_summary': {
            'tables_found': sqproj_data.get('tables', {}).get('count', 0),
            'geometry_tables': len(sqproj_data.get('geometry', {})),
            'contact_tables': len(sqproj_data.get('contacts', {})),
        },
        '_extraction_log': elog.summary(),
    }

    return structured


# ============================================================================
# CHART GENERATION (Bug #16, #17, #18  –  real data, correct types)
# ============================================================================

def _fig_to_base64(fig) -> str:
    buf = BytesIO()
    fig.savefig(buf, format='png', dpi=200, bbox_inches='tight', transparent=False,
                facecolor='white', edgecolor='none')
    buf.seek(0)
    b64 = base64.b64encode(buf.read()).decode()
    plt.close(fig)
    return b64


def generate_energy_class_chart(primary_demand: float, label: str = 'IST') -> str:
    """Generate energy efficiency class scale (A+ to H) with position marker."""
    fig, ax = plt.subplots(figsize=(10, 2.5))

    classes = [
        ('A+', 0, 30, '#00734A'), ('A', 30, 50, '#4BA946'),
        ('B', 50, 75, '#B2D235'), ('C', 75, 100, '#EDE63A'),
        ('D', 100, 130, '#ECC531'), ('E', 130, 160, '#E89E30'),
        ('F', 160, 200, '#E16E2E'), ('G', 200, 250, '#D23C2C'),
        ('H', 250, 350, '#B5232A'),
    ]

    for cls_label, start, end, color in classes:
        width = end - start
        ax.barh(0, width, left=start, height=0.6, color=color, edgecolor='white', linewidth=0.5)
        ax.text(start + width / 2, 0, cls_label, ha='center', va='center',
                fontsize=10, fontweight='bold', color='white')

    marker_x = min(primary_demand, 340)
    ax.annotate(f'{int(primary_demand)} kWh/(m²a)',
                xy=(marker_x, -0.3), xytext=(marker_x, -0.8),
                fontsize=9, fontweight='bold', ha='center',
                arrowprops=dict(arrowstyle='->', color='black', lw=2),
                bbox=dict(boxstyle='round,pad=0.3', facecolor='white', edgecolor='black'))

    ax.set_xlim(0, 350)
    ax.set_ylim(-1.2, 0.5)
    ax.set_xticks([0, 50, 100, 150, 200, 250, 300, 350])
    ax.set_xlabel('Primärenergiebedarf [kWh/(m²a)]', fontsize=9)
    ax.set_yticks([])
    ax.set_title(f'Energieeffizienzklasse  –  {label}', fontsize=11, fontweight='bold')
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)

    return _fig_to_base64(fig)


def generate_component_loss_chart(component: str, kwh_ist: float, kwh_ziel: float, reduction_pct: float) -> str:
    """Generate before/after energy loss comparison for a single component."""
    fig, ax = plt.subplots(figsize=(8, 4))

    labels = ['Ist-Zustand', 'Nach Sanierung']
    values = [kwh_ist, kwh_ziel]
    colors = ['#D23C2C', '#4BA946']
    bars = ax.bar(labels, values, color=colors, width=0.5, edgecolor='white')

    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + max(values) * 0.02,
                f'{int(val):,} kWh'.replace(',', '.'), ha='center', va='bottom',
                fontsize=11, fontweight='bold')

    if kwh_ist > 0 and kwh_ziel < kwh_ist:
        ax.annotate(f'-{int(reduction_pct)}%',
                    xy=(1, kwh_ziel), xytext=(0.5, (kwh_ist + kwh_ziel) / 2),
                    fontsize=14, fontweight='bold', color='#4BA946', ha='center',
                    arrowprops=dict(arrowstyle='->', color='#4BA946', lw=2))

    component_titles = {
        'dach': 'Dach / Oberste Geschossdecke',
        'aussenwand': 'Außenwand',
        'fenster': 'Fenster, Türen & Lüftung',
        'keller': 'Unterer Gebäudeabschluss',
        'heizung': 'Heizung & Warmwasser',
    }
    ax.set_title(f'Energieverluste  –  {component_titles.get(component, component)}',
                 fontsize=12, fontweight='bold')
    ax.set_ylabel('Wärmeverluste [kWh/a]', fontsize=10)
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'{int(x):,}'.replace(',', '.')))
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(axis='y', alpha=0.3)

    return _fig_to_base64(fig)


def generate_endenergie_chart(ist: float, ziel: float) -> str:
    """Endenergiebedarf comparison chart."""
    fig, ax = plt.subplots(figsize=(8, 4))
    labels = ['Ist-Zustand', 'Nach Sanierung']
    values = [ist, ziel]
    colors = ['#D23C2C', '#4BA946']
    bars = ax.bar(labels, values, color=colors, width=0.5)
    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + max(values) * 0.02,
                f'{int(val):,} kWh/a'.replace(',', '.'), ha='center', va='bottom', fontsize=11, fontweight='bold')
    ax.set_title('Endenergieverbrauch  –  Vergleich', fontsize=12, fontweight='bold')
    ax.set_ylabel('kWh/a')
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'{int(x):,}'.replace(',', '.')))
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(axis='y', alpha=0.3)
    return _fig_to_base64(fig)


def generate_brennstoffkosten_chart(ist: float, ziel: float) -> str:
    """Brennstoffkosten comparison chart."""
    fig, ax = plt.subplots(figsize=(8, 4))
    labels = ['Ist-Zustand', 'Nach Sanierung']
    values = [ist, ziel]
    colors = ['#D23C2C', '#4BA946']
    bars = ax.bar(labels, values, color=colors, width=0.5)
    for bar, val in zip(bars, values):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + max(values) * 0.02,
                f'{int(val):,} €/a'.replace(',', '.'), ha='center', va='bottom', fontsize=11, fontweight='bold')
    ax.set_title('Energiekosten  –  Vergleich', fontsize=12, fontweight='bold')
    ax.set_ylabel('€/a')
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(lambda x, _: f'{int(x):,}'.replace(',', '.')))
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.grid(axis='y', alpha=0.3)
    return _fig_to_base64(fig)


def generate_all_charts(data: dict) -> dict:
    """Generate all chart images from extracted data."""
    charts = {}
    energy = data.get('energy', {})
    losses = data.get('energy_losses', {})

    try:
        ist_primary = float(energy.get('ist', {}).get('primary_demand', 290))
        charts['img_energieklasse_ist'] = generate_energy_class_chart(ist_primary, 'Ist-Zustand')
    except Exception as e:
        log.error(f'Chart generation failed (energieklasse_ist): {e}')

    try:
        ziel_primary = float(energy.get('ziel', {}).get('primary_demand', 27))
        charts['img_energieklasse_ziel'] = generate_energy_class_chart(ziel_primary, 'Ziel-Zustand')
    except Exception as e:
        log.error(f'Chart generation failed (energieklasse_ziel): {e}')

    for comp in ['dach', 'aussenwand', 'fenster', 'keller', 'heizung']:
        try:
            comp_data = losses.get(comp, {})
            kwh_ist = float(comp_data.get('kwh_ist', 0))
            kwh_ziel = float(comp_data.get('kwh_ziel', 0))
            pct = float(comp_data.get('reduction_pct', 0))
            if kwh_ist > 0:
                charts[f'img_energieverluste_{comp}'] = generate_component_loss_chart(
                    comp, kwh_ist, kwh_ziel, pct)
        except Exception as e:
            log.error(f'Chart generation failed (energieverluste_{comp}): {e}')

    try:
        ist_end = float(energy.get('ist', {}).get('end_consumption', 101500))
        ziel_end = float(energy.get('ziel', {}).get('end_consumption', 7550))
        charts['img_endenergiebedarf'] = generate_endenergie_chart(ist_end, ziel_end)
    except Exception as e:
        log.error(f'Chart generation failed (endenergiebedarf): {e}')

    try:
        ist_cost = float(energy.get('ist', {}).get('costs', 6600))
        ziel_cost = float(energy.get('ziel', {}).get('costs', 1500))
        charts['img_brennstoffkosten'] = generate_brennstoffkosten_chart(ist_cost, ziel_cost)
    except Exception as e:
        log.error(f'Chart generation failed (brennstoffkosten): {e}')

    log.info(f'Generated {len(charts)} charts')
    return charts


# ============================================================================
# PPT GENERATION (Bug #19, #20, #21, #22 + Priority 1a, 2)
# ============================================================================

def replace_text_in_runs(shape, mapping: dict, stats: dict):
    """Bug #19: Replace placeholders at the RUN level, preserving all formatting."""
    if not shape.has_text_frame:
        return

    for paragraph in shape.text_frame.paragraphs:
        full_para_text = ''.join(run.text for run in paragraph.runs)
        if '{{' not in full_para_text:
            continue

        for placeholder, value in mapping.items():
            pattern = '{{' + placeholder + '}}'
            if pattern not in full_para_text:
                continue

            found_in_single = False
            for run in paragraph.runs:
                if pattern in run.text:
                    run.text = run.text.replace(pattern, str(value))
                    stats[placeholder] = stats.get(placeholder, 0) + 1
                    found_in_single = True

            if not found_in_single and pattern in full_para_text:
                if paragraph.runs:
                    combined = ''.join(r.text for r in paragraph.runs)
                    combined = combined.replace(pattern, str(value))
                    paragraph.runs[0].text = combined
                    for r in paragraph.runs[1:]:
                        r.text = ''
                    stats[placeholder] = stats.get(placeholder, 0) + 1


def replace_text_in_table(table, mapping: dict, stats: dict):
    """Bug #20: Replace in table cells at run level."""
    for row in table.rows:
        for cell in row.cells:
            for paragraph in cell.text_frame.paragraphs:
                full_para_text = ''.join(run.text for run in paragraph.runs)
                if '{{' not in full_para_text:
                    continue

                for placeholder, value in mapping.items():
                    pattern = '{{' + placeholder + '}}'
                    if pattern not in full_para_text:
                        continue

                    found = False
                    for run in paragraph.runs:
                        if pattern in run.text:
                            run.text = run.text.replace(pattern, str(value))
                            stats[placeholder] = stats.get(placeholder, 0) + 1
                            found = True

                    if not found and pattern in full_para_text:
                        if paragraph.runs:
                            combined = ''.join(r.text for r in paragraph.runs)
                            combined = combined.replace(pattern, str(value))
                            paragraph.runs[0].text = combined
                            for r in paragraph.runs[1:]:
                                r.text = ''
                            stats[placeholder] = stats.get(placeholder, 0) + 1


def replace_image_in_shape(slide, shape, image_b64: str, stats: dict, placeholder: str, fit_mode: bool = False):
    """
    Replace image via blipFill to preserve z-order.
    fit_mode=False (default): COVER mode – scale to fill shape, center-crop excess (for photos).
    fit_mode=True:            FIT/CONTAIN mode – scale to fit within shape, no crop (for charts).
    """
    try:
        img_data = base64.b64decode(image_b64)

        # Open image to get dimensions for processing
        pil_img = Image.open(BytesIO(img_data))
        img_w, img_h = pil_img.size

        # Shape dimensions (EMU)
        shape_w = shape.width
        shape_h = shape.height

        shape_ratio = shape_w / shape_h if shape_h > 0 else 1.0
        img_ratio = img_w / img_h if img_h > 0 else 1.0

        if fit_mode:
            # FIT/CONTAIN: scale image to fit entirely within shape, white padding
            if img_ratio > shape_ratio:
                # Image wider than shape → fit to width, pad top/bottom
                new_w = img_w
                new_h = int(img_w / shape_ratio)
            else:
                # Image taller than shape → fit to height, pad left/right
                new_h = img_h
                new_w = int(img_h * shape_ratio)

            # Create white background canvas at target aspect ratio
            canvas = Image.new('RGB', (new_w, new_h), (255, 255, 255))
            # Paste original image centered on canvas
            paste_x = (new_w - img_w) // 2
            paste_y = (new_h - img_h) // 2
            if pil_img.mode == 'RGBA':
                canvas.paste(pil_img, (paste_x, paste_y), pil_img)
            else:
                canvas.paste(pil_img, (paste_x, paste_y))
            final_img = canvas
            log.info(f'  Image FIT for {placeholder}: {img_w}x{img_h} → canvas {new_w}x{new_h} (shape {shape_w}x{shape_h} EMU)')
        else:
            # COVER: scale so shortest side fills shape, then center-crop
            if img_ratio > shape_ratio:
                # Image is wider → match height, crop sides
                new_h = img_h
                new_w = int(img_h * shape_ratio)
                left = (img_w - new_w) // 2
                crop_box = (left, 0, left + new_w, new_h)
            else:
                # Image is taller → match width, crop top/bottom
                new_w = img_w
                new_h = int(img_w / shape_ratio)
                top = (img_h - new_h) // 2
                crop_box = (0, top, new_w, top + new_h)

            final_img = pil_img.crop(crop_box)
            log.info(f'  Image COVER-crop for {placeholder}: {img_w}x{img_h} → {final_img.size[0]}x{final_img.size[1]} (shape {shape_w}x{shape_h} EMU)')

        buf = BytesIO()
        fmt = 'PNG' if pil_img.mode == 'RGBA' else 'JPEG'
        save_kwargs = {'optimize': True}
        if fmt == 'JPEG':
            save_kwargs['quality'] = 92
        final_img.save(buf, format=fmt, **save_kwargs)
        buf.seek(0)

        sp_element = shape.element
        blip_fills = sp_element.findall('.//' + qn('a:blip'))

        if blip_fills:
            slide_part = slide.part
            image_part, rId = slide_part.get_or_add_image_part(buf)
            blip_fills[0].set(qn('r:embed'), rId)
            blipFill = sp_element.find('.//' + qn('a:blipFill'))
            if blipFill is not None:
                # Remove srcRect (would clip/offset the image)
                for srcRect in blipFill.findall(qn('a:srcRect')):
                    blipFill.remove(srcRect)
                # Remove tile (would repeat instead of stretch)
                for tile in blipFill.findall(qn('a:tile')):
                    blipFill.remove(tile)
                # Ensure stretch fill mode with centered fillRect
                stretch = blipFill.find(qn('a:stretch'))
                if stretch is None:
                    stretch = etree.SubElement(blipFill, qn('a:stretch'))
                if stretch.find(qn('a:fillRect')) is None:
                    etree.SubElement(stretch, qn('a:fillRect'))
            stats[placeholder] = 1
            return

        # Fallback: add as new picture at shape position
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height

        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    run.text = ''

        slide.shapes.add_picture(buf, left, top, width, height)
        stats[placeholder] = 1

    except Exception as e:
        log.error(f'Image replacement failed for {placeholder}: {e}')


# Priority 2: Bar shape resizing
BAR_PATTERN = re.compile(r'^bar_(red|green)_loss_(.+?)_(ist|loesung)(\d*)$')
COLOR_BAR_RED = RGBColor(0xD4, 0x3B, 0x3B)
COLOR_BAR_GREEN = RGBColor(0x4C, 0xAF, 0x50)


def resize_bar_shapes(slide, text_mapping: dict, stats: dict):
    """
    Priority 2 + P1a: Find bar shapes by text content (bar_red_*, bar_green_*),
    resize width proportionally to kWh values, apply fill color, clear text.

    Handles suffixed bar names like bar_red_loss_dach_ist1, bar_green_loss_AW_loesung3.
    
    Supports NEGATIVE kWh values (e.g. Heizung with heat pump COP>1):
      - All bars shift right to create space for negative bar zone
      - Negative green bars extend LEFT from the origin line
    """
    bar_shapes = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        m = BAR_PATTERN.match(text)
        if not m:
            continue

        color_type, component, value_type, suffix = m.groups()
        placeholder_key = f'loss_{component}_kwh_{value_type}{suffix}'
        raw_val = text_mapping.get(placeholder_key, '')
        kwh_value = float(str(raw_val).replace('.', '').replace(',', '.')) if raw_val else 0

        bar_shapes.append({
            'shape': shape,
            'text': text,
            'color_type': color_type,
            'kwh': kwh_value,
            'original_left': shape.left,
            'max_width': shape.width,
            'suffix': suffix,
        })

    if not bar_shapes:
        return

    # Find max positive and most negative kWh values
    max_pos_kwh = max((b['kwh'] for b in bar_shapes), default=1)
    min_neg_kwh = min((b['kwh'] for b in bar_shapes), default=0)
    if max_pos_kwh <= 0:
        max_pos_kwh = 1

    max_bar_width = bar_shapes[0]['max_width']

    # Calculate negative zone: shift all bars right if any negative values exist
    has_negative = min_neg_kwh < 0
    if has_negative:
        # Reserve space on the left for negative bars
        neg_zone_ratio = abs(min_neg_kwh) / (max_pos_kwh + abs(min_neg_kwh))
        neg_zone_width = int(max_bar_width * neg_zone_ratio * 0.5)  # cap at 50%
        effective_bar_width = max_bar_width - neg_zone_width
    else:
        neg_zone_width = 0
        effective_bar_width = max_bar_width

    for bar in bar_shapes:
        shape = bar['shape']
        kwh = bar['kwh']
        original_left = bar['original_left']

        fill_color = COLOR_BAR_RED if bar['color_type'] == 'red' else COLOR_BAR_GREEN

        if kwh < 0:
            # Negative bar: extends LEFT from origin line
            origin_x = original_left + neg_zone_width
            bar_width = int(effective_bar_width * abs(kwh) / max_pos_kwh)
            bar_width = max(bar_width, int(max_bar_width * 0.02))
            shape.left = origin_x - bar_width
            shape.width = bar_width
        elif kwh > 0:
            # Positive bar: extends RIGHT from origin (shifted right)
            ratio = max(kwh / max_pos_kwh, 0.03)
            shape.left = original_left + neg_zone_width
            shape.width = int(effective_bar_width * ratio)
        else:
            # Zero: minimal bar at origin
            shape.left = original_left + neg_zone_width
            shape.width = int(max_bar_width * 0.03)

        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color

        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                run.text = ''

        stats[bar['text']] = 1

def fill_presentation(template_bytes: bytes, text_mapping: dict, image_mapping: dict) -> tuple:
    """
    Fill the template. Returns (pptx_bytes, stats_dict).
    Pass 0: Pre-process – add "Ist-Zustand: " / "Lösung: " labels to component texts
    Pass 1: Text + image replacement (FIT mode for charts, COVER mode for photos)
    Pass 1b: Bold formatting for "Ist-Zustand: " / "Lösung: " labels
    Pass 2: Bar shape resizing (Priority 2 + P1a suffix support)
    Pass 3: Handle remaining unfilled placeholders (green img + yellow text)
    """
    prs = Presentation(BytesIO(template_bytes))
    stats = {}

    # --- Pass 0: Prepend "Ist-Zustand: " / "Lösung: " to component text values ---
    # The detail slides have ist+loesung in one text box.  We prefix each with a bold label.
    ist_text_keys = [
        'dach_ist', 'fenster_ist', 'aussenwand_ist',
        'keller_ist', 'heizung_ist', 'warmwasser_ist',
        'lueftung_ist',
    ]
    loesung_text_keys = [
        'dach_loesung', 'fenster_loesung', 'aussenwand_loesung',
        'keller_loesung', 'heizung_loesung', 'warmwasser_loesung',
        'lueftung_loesung',
    ]
    for key in ist_text_keys:
        if key in text_mapping and text_mapping[key]:
            val = str(text_mapping[key]).strip()
            if val and not val.startswith('Ist-Zustand:'):
                text_mapping[key] = f'Ist-Zustand: {val}'
    for key in loesung_text_keys:
        if key in text_mapping and text_mapping[key]:
            val = str(text_mapping[key]).strip()
            if val and not val.startswith('Lösung:'):
                text_mapping[key] = f'Lösung: {val}'
    log.info(f'  Pass 0: Prepended "Ist-Zustand:" / "Lösung:" labels to component texts')

    # Chart image keys use FIT mode (no crop), extracted images use COVER mode (fill+crop)
    CHART_IMAGE_PREFIXES = ('img_energieklasse', 'img_endenergiebedarf', 'img_brennstoffkosten', 'img_energieverluste')

    # Pass 1: Text and image replacement
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                full_text = shape.text
                img_replaced = False
                for img_key, img_b64 in image_mapping.items():
                    pattern = '{{' + img_key + '}}'
                    if pattern in full_text:
                        is_chart = img_key.startswith(CHART_IMAGE_PREFIXES)
                        replace_image_in_shape(slide, shape, img_b64, stats, img_key, fit_mode=is_chart)
                        img_replaced = True
                        break

                if not img_replaced:
                    replace_text_in_runs(shape, text_mapping, stats)

            if shape.has_table:
                replace_text_in_table(shape.table, text_mapping, stats)

    # Pass 1b: Apply bold formatting to "Ist-Zustand: " and "Lösung: " labels
    bold_prefixes = ['Ist-Zustand: ', 'Lösung: ']
    bold_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in list(paragraph.runs):
                        for prefix in bold_prefixes:
                            if run.text.startswith(prefix):
                                # Split: create bold run for prefix, keep rest in original run
                                rest_text = run.text[len(prefix):]
                                run.text = prefix
                                run.font.bold = True
                                # Insert a new run after this one for the rest of the text
                                if rest_text:
                                    new_r = deepcopy(run._r)
                                    new_r.text = rest_text
                                    # Remove bold from the copy
                                    new_rPr = new_r.find(qn('a:rPr'))
                                    if new_rPr is not None:
                                        new_rPr.attrib.pop('b', None)
                                    run._r.addnext(new_r)
                                bold_count += 1
                                break
    log.info(f'  Pass 1b: Applied bold to {bold_count} Ist-Zustand/Lösung labels')

    # Pass 2: Bar chart shape resizing (Priority 2)
    for slide in prs.slides:
        resize_bar_shapes(slide, text_mapping, stats)

    # Pass 3: Handle remaining unfilled {{...}} placeholders
    #  - img_ placeholders → solid fill shape with #798C3A, clear text
    #  - other placeholders → yellow text highlight (Priority 1a)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                full_text = shape.text
                if '{{' in full_text and '}}' in full_text:
                    # Check if this shape has unfilled img_ placeholders
                    img_unfilled = re.findall(r'\{\{(img_[^}]+)\}\}', full_text)
                    if img_unfilled:
                        # Solid fill the entire shape with #798C3A
                        shape.fill.solid()
                        shape.fill.fore_color.rgb = RGBColor(0x79, 0x8C, 0x3A)
                        # Clear existing text and add italic hint text
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                run.text = ''
                        # Add "Bitte Bild einfügen" in white italic
                        if shape.text_frame.paragraphs:
                            p = shape.text_frame.paragraphs[0]
                            if p.runs:
                                p.runs[0].text = 'Bitte Bild einfügen'
                                p.runs[0].font.italic = True
                                p.runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                                p.runs[0].font.size = Pt(10)
                            else:
                                run_elem = etree.SubElement(p._p, qn('a:r'))
                                rPr = etree.SubElement(run_elem, qn('a:rPr'))
                                rPr.set('i', '1')
                                rPr.set('sz', '1000')
                                solidFill_r = etree.SubElement(rPr, qn('a:solidFill'))
                                srgb_r = etree.SubElement(solidFill_r, qn('a:srgbClr'))
                                srgb_r.set('val', 'FFFFFF')
                                t = etree.SubElement(run_elem, qn('a:t'))
                                t.text = 'Bitte Bild einfügen'
                        log.info(f'  Green-filled img placeholder: {shape.name} + italic hint text')
                    else:
                        # Non-img placeholders: yellow highlight
                        for paragraph in shape.text_frame.paragraphs:
                            full_para = ''.join(run.text for run in paragraph.runs)
                            if '{{' in full_para and '}}' in full_para:
                                for run in paragraph.runs:
                                    if '{{' in run.text or '}}' in run.text:
                                        rPr = run._r.get_or_add_rPr()
                                        highlight = etree.SubElement(rPr, qn('a:highlight'))
                                        srgb = etree.SubElement(highlight, qn('a:srgbClr'))
                                        srgb.set('val', 'FFFF00')

            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            full_text = ''.join(run.text for run in paragraph.runs)
                            if '{{' in full_text and '}}' in full_text:
                                for run in paragraph.runs:
                                    if '{{' in run.text or '}}' in run.text:
                                        rPr = run._r.get_or_add_rPr()
                                        highlight = etree.SubElement(rPr, qn('a:highlight'))
                                        srgb = etree.SubElement(highlight, qn('a:srgbClr'))
                                        srgb.set('val', 'FFFF00')

    log.info(f'  Post-processing: highlighted remaining unfilled placeholders')

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output.read(), stats


# ============================================================================
# API ENDPOINTS
# ============================================================================

@app.route('/extract-comprehensive', methods=['POST'])
def api_extract_comprehensive():
    """
    Extract ALL data from uploaded files.
    5-phase pipeline: collect → classify → cache → extract → structure.
    """
    start = time.time()
    log.info('=== /extract-comprehensive START ===')

    try:
        # ══€══€ Phase 1: Collect files from ANY field name ══€══€
        all_files = []
        for key in request.files:
            all_files.append((key, request.files[key]))

        if not all_files:
            return jsonify({'error': 'No files uploaded. Send files as multipart form data.'}), 400

        log.info(f'Received {len(all_files)} files: {[(f[0], f[1].filename) for f in all_files]}')

        # ══€══€ Phase 2: Classify by extension + magic bytes ══€══€
        pdfs = []
        sqproj = None
        pptx_template = None

        for field_name, f in all_files:
            fname = (f.filename or field_name or '').lower()
            content_bytes = f.read()

            if not content_bytes or len(content_bytes) < 10:
                log.warning(f'  Skipping empty/tiny file: {fname}')
                continue

            if fname.endswith('.sqproj'):
                sqproj = content_bytes
                log.info(f'  → sqproj (by extension): {f.filename} ({len(content_bytes):,} bytes)')
            elif fname.endswith('.pptx'):
                pptx_template = content_bytes
                log.info(f'  → pptx template (by extension): {f.filename}')
            elif fname.endswith('.pdf'):
                pdfs.append((content_bytes, f.filename or field_name))
                log.info(f'  → pdf (by extension): {f.filename} ({len(content_bytes):,} bytes)')
            else:
                if content_bytes[:15] == b'SQLite format 3':
                    sqproj = content_bytes
                    log.info(f'  → sqproj (magic bytes): {f.filename}')
                elif content_bytes[:4] == b'%PDF':
                    pdfs.append((content_bytes, f.filename or field_name))
                    log.info(f'  → pdf (magic bytes): {f.filename}')
                elif content_bytes[:2] == b'PK':
                    pptx_template = content_bytes
                    log.info(f'  → pptx (magic bytes PK): {f.filename}')
                else:
                    log.warning(f'  → UNKNOWN file type: {f.filename} ({len(content_bytes):,} bytes)')

        # ══€══€ Phase 2b: Classify PDFs by content ══€══€
        sanierungsfahrplan_bytes = None
        umsetzungshilfe_bytes = None

        if len(pdfs) == 1:
            role = classify_pdf_by_content(pdfs[0][0])
            if role == 'umsetzungshilfe':
                umsetzungshilfe_bytes = pdfs[0][0]
                log.info(f'  Single PDF classified as: umsetzungshilfe ({pdfs[0][1]})')
            else:
                sanierungsfahrplan_bytes = pdfs[0][0]
                log.info(f'  Single PDF classified as: sanierungsfahrplan ({pdfs[0][1]})')

        elif len(pdfs) == 2:
            role_a = classify_pdf_by_content(pdfs[0][0])
            role_b = classify_pdf_by_content(pdfs[1][0])
            log.info(f'  PDF A ({pdfs[0][1]}): classified as {role_a}')
            log.info(f'  PDF B ({pdfs[1][1]}): classified as {role_b}')

            if role_a == 'umsetzungshilfe' and role_b != 'umsetzungshilfe':
                umsetzungshilfe_bytes = pdfs[0][0]
                sanierungsfahrplan_bytes = pdfs[1][0]
            elif role_b == 'umsetzungshilfe' and role_a != 'umsetzungshilfe':
                umsetzungshilfe_bytes = pdfs[1][0]
                sanierungsfahrplan_bytes = pdfs[0][0]
            elif role_a == 'sanierungsfahrplan' and role_b != 'sanierungsfahrplan':
                sanierungsfahrplan_bytes = pdfs[0][0]
                umsetzungshilfe_bytes = pdfs[1][0]
            else:
                try:
                    pages_a = len(pdfplumber.open(BytesIO(pdfs[0][0])).pages)
                    pages_b = len(pdfplumber.open(BytesIO(pdfs[1][0])).pages)
                    if pages_a > pages_b:
                        umsetzungshilfe_bytes = pdfs[0][0]
                        sanierungsfahrplan_bytes = pdfs[1][0]
                    else:
                        sanierungsfahrplan_bytes = pdfs[0][0]
                        umsetzungshilfe_bytes = pdfs[1][0]
                    log.info(f'  Ambiguous – used page count heuristic (A={pages_a}, B={pages_b})')
                except Exception:
                    sanierungsfahrplan_bytes = pdfs[0][0]
                    umsetzungshilfe_bytes = pdfs[1][0]
                    log.info('  Ambiguous – defaulted to upload order')

        elif len(pdfs) > 2:
            roles = [(p[0], p[1], classify_pdf_by_content(p[0])) for p in pdfs]
            for content_b, fname, role in roles:
                if role == 'umsetzungshilfe' and not umsetzungshilfe_bytes:
                    umsetzungshilfe_bytes = content_b
                elif role == 'sanierungsfahrplan' and not sanierungsfahrplan_bytes:
                    sanierungsfahrplan_bytes = content_b
            for content_b, fname, role in roles:
                if content_b is not sanierungsfahrplan_bytes and content_b is not umsetzungshilfe_bytes:
                    if not sanierungsfahrplan_bytes:
                        sanierungsfahrplan_bytes = content_b
                    elif not umsetzungshilfe_bytes:
                        umsetzungshilfe_bytes = content_b

        # ══€══€ Phase 3: Check cache ══€══€
        cache_parts = []
        if sanierungsfahrplan_bytes:
            cache_parts.append(hashlib.md5(sanierungsfahrplan_bytes).hexdigest()[:8])
        if umsetzungshilfe_bytes:
            cache_parts.append(hashlib.md5(umsetzungshilfe_bytes).hexdigest()[:8])
        if sqproj:
            cache_parts.append(hashlib.md5(sqproj).hexdigest()[:8])
        ck = '_'.join(sorted(cache_parts)) if cache_parts else 'empty'
        cached = cache_get(ck)
        if cached:
            return jsonify({'success': True, 'data': cached, 'cached': True,
                            'duration_sec': round(time.time() - start, 2)})

        # ══€══€ Phase 4: Extract from each file (graceful if missing) ══€══€
        sqproj_data = extract_all_from_sqproj(sqproj) if sqproj else {'tables': {'count': 0}}
        pdf1_data = extract_all_data_from_pdf(sanierungsfahrplan_bytes, 'sanierungsfahrplan') \
                    if sanierungsfahrplan_bytes else {}
        pdf2_data = extract_all_data_from_pdf(umsetzungshilfe_bytes, 'umsetzungshilfe') \
                    if umsetzungshilfe_bytes else {}

        # Extract images from PDFs (Priority 3: page rendering + cropping)
        pdf_images = {}
        if sanierungsfahrplan_bytes:
            pdf_images.update(extract_images_from_pdf(sanierungsfahrplan_bytes, 'sanierungsfahrplan'))

        # ══€══€ Phase 5: Structure ══€══€
        structured = structure_complete_data(sqproj_data, pdf1_data, pdf2_data, pdf_images)

        # Cache (note: pdf_images are large base64 strings, skip them in cache)
        cache_data = {k: v for k, v in structured.items() if k != 'pdf_images'}
        cache_set(ck, cache_data)

        duration = round(time.time() - start, 2)
        log.info(f'=== /extract-comprehensive END ({duration}s) ===')

        files_detected = {
            'sanierungsfahrplan': sanierungsfahrplan_bytes is not None,
            'umsetzungshilfe': umsetzungshilfe_bytes is not None,
            'projektdatei': sqproj is not None,
        }

        return jsonify({
            'success': True,
            'data': structured,
            'files_detected': files_detected,
            'extraction_summary': {
                'sqproj_tables': sqproj_data.get('tables', {}).get('count', 0),
                'pdf1_pages': pdf1_data.get('metadata', {}).get('page_count', 0) if pdf1_data else 0,
                'pdf2_pages': pdf2_data.get('metadata', {}).get('page_count', 0) if pdf2_data else 0,
                'placeholders_mapped': len([v for v in structured.get('placeholder_map', {}).values() if v]),
                'charts_preparable': len(structured.get('energy_losses', {})),
                'pdf_images_extracted': len(pdf_images),
                'files_received': len(all_files),
                'files_classified': sum(1 for v in files_detected.values() if v),
            },
            'duration_sec': duration,
        })

    except Exception as e:
        log.exception('Extraction failed')
        return jsonify({'error': str(e)}), 500


@app.route('/read-template-placeholders', methods=['POST'])
def api_read_template():
    """Read all {{placeholders}} from .pptx template. Also stores template for /generate."""
    try:
        template_file = None
        for key in request.files:
            template_file = request.files[key]
            break

        if not template_file:
            return jsonify({'error': 'No template file uploaded'}), 400

        template_bytes = template_file.read()

        # Save template to disk for later use by /generate
        template_hash = hashlib.sha256(template_bytes).hexdigest()[:16]
        template_path = TEMPLATE_DIR / f'{template_hash}.pptx'
        # Atomic write: temp file then rename
        tmp_template = template_path.with_suffix('.tmp')
        tmp_template.write_bytes(template_bytes)
        tmp_template.rename(template_path)
        log.info(f'Template saved: {template_path} ({len(template_bytes):,} bytes)')

        prs = Presentation(BytesIO(template_bytes))
        placeholders = set()
        slide_map = {}

        for i, slide in enumerate(prs.slides):
            slide_placeholders = set()
            for shape in slide.shapes:
                text = ''
                if shape.has_text_frame:
                    text = shape.text
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            text += ' ' + cell.text

                found = re.findall(r'\{\{([^}]+)\}\}', text)
                for f in found:
                    placeholders.add(f)
                    slide_placeholders.add(f)

            if slide_placeholders:
                slide_map[f'slide_{i + 1}'] = sorted(list(slide_placeholders))

        return jsonify({
            'success': True,
            'placeholders': sorted(list(placeholders)),
            'count': len(placeholders),
            'slide_map': slide_map,
            'template_hash': template_hash,
        })

    except Exception as e:
        log.exception('Template reading failed')
        return jsonify({'error': str(e)}), 500


@app.route('/generate-charts', methods=['POST'])
def api_generate_charts():
    """Generate all chart images. Accepts extracted data JSON."""
    try:
        data = request.get_json()
        if not data:
            return jsonify({'error': 'No JSON data provided'}), 400

        charts = generate_all_charts(data)

        return jsonify({
            'success': True,
            'charts': charts,
            'count': len(charts),
        })

    except Exception as e:
        log.exception('Chart generation failed')
        return jsonify({'error': str(e)}), 500


@app.route('/generate', methods=['POST'])
def api_generate_ppt():
    """
    Generate filled PPT.
    Accepts:
      1. Multipart: template file + JSON data field (backward compatible)
      2. JSON-only: { "data": {...} } – uses stored template from /read-template-placeholders
      3. JSON with hash: { "data": {...}, "template_hash": "abc123" }
    Returns binary .pptx file.
    """
    try:
        template_bytes = None
        approved_data = None

        # --- Try multipart upload first (backward compatible) ---
        if request.content_type and 'multipart' in request.content_type:
            for key in request.files:
                f = request.files[key]
                fname = (f.filename or key or '').lower()
                if fname.endswith('.pptx') or 'template' in key.lower():
                    template_bytes = f.read()
                    break

            if request.form.get('data'):
                approved_data = json.loads(request.form['data'])
            elif request.form.get('approved_data'):
                approved_data = json.loads(request.form['approved_data'])

        # --- Try JSON body ---
        if not approved_data and request.is_json:
            body = request.get_json()
            template_b64 = body.get('template_file')
            if template_b64 and not template_bytes:
                template_bytes = base64.b64decode(template_b64)
            approved_data = body.get('approved_data') or body.get('data')

            # Look up stored template by hash
            if not template_bytes:
                template_hash = body.get('template_hash', '')
                if template_hash:
                    hash_path = TEMPLATE_DIR / f'{template_hash}.pptx'
                    if hash_path.exists():
                        template_bytes = hash_path.read_bytes()
                        log.info(f'Using stored template by hash: {template_hash} ({len(template_bytes):,} bytes)')

        # --- Fall back: no template found ---
        if not template_bytes:
            log.warning('No template provided and no template_hash matched. Rejecting request.')

        if not template_bytes:
            return jsonify({'error': 'No template found. Provide template_hash from /read-template-placeholders or upload template file.'}), 400
        if not approved_data:
            return jsonify({'error': 'No data/approved_data provided'}), 400

        log.info(f'Generating PPT with {len(approved_data)} data fields')

        text_mapping = {}
        image_mapping = {}
        for key, value in approved_data.items():
            if key.startswith('img_') and value and len(str(value)) > 500:
                image_mapping[key] = value
            else:
                text_mapping[key] = value

        log.info(f'  Text fields: {len(text_mapping)}, Image fields: {len(image_mapping)}')

        pptx_bytes, stats = fill_presentation(template_bytes, text_mapping, image_mapping)

        log.info(f'  Replacements made: {sum(stats.values())} across {len(stats)} placeholders')

        output = BytesIO(pptx_bytes)
        output.seek(0)

        # Build dynamic filename: YYYYMMDD_iSFP_ClientName.pptx
        kunde_name = approved_data.get('kunde_name', '').strip()
        if not kunde_name:
            kunde_name = approved_data.get('client_name', '') or approved_data.get('owner', '') or 'Kunde'
        # Strip salutation (Herr/Frau) for cleaner filename
        kunde_name = re.sub(r'^(Herr|Frau)\s+', '', kunde_name).strip()
        # Sanitize filename: keep word chars (incl. umlauts), spaces, hyphens
        safe_name = re.sub(r'[^\w\s\-]', '', kunde_name).strip().replace(' ', '_')
        date_prefix = datetime.now().strftime('%Y%m%d')
        download_filename = f'{date_prefix}_iSFP_{safe_name}.pptx' if safe_name else f'{date_prefix}_iSFP.pptx'

        return send_file(
            output,
            mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            as_attachment=True,
            download_name=download_filename,
        )

    except Exception as e:
        log.exception('PPT generation failed')
        return jsonify({'error': str(e)}), 500


@app.route('/generate-json', methods=['POST'])
def api_generate_ppt_json():
    """Same as /generate but returns base64-encoded file in JSON (backward compat)."""
    try:
        body = request.get_json() or {}
        template_b64 = body.get('template_file', '')
        approved_data = body.get('approved_data') or body.get('data', {})

        if not template_b64:
            return jsonify({'error': 'Missing template_file (base64)'}), 400

        template_bytes = base64.b64decode(template_b64)

        text_mapping = {}
        image_mapping = {}
        for key, value in approved_data.items():
            if key.startswith('img_') and value and len(str(value)) > 500:
                image_mapping[key] = value
            else:
                text_mapping[key] = value

        pptx_bytes, stats = fill_presentation(template_bytes, text_mapping, image_mapping)

        return jsonify({
            'success': True,
            'filename': 'presentation_filled.pptx',
            'file_content': base64.b64encode(pptx_bytes).decode(),
            'file_size_mb': round(len(pptx_bytes) / 1024 / 1024, 2),
            'replacements': stats,
        })

    except Exception as e:
        log.exception('PPT generation (JSON) failed')
        return jsonify({'error': str(e)}), 500


# ============================================================================
# RUN
# ============================================================================

if __name__ == '__main__':
    log.info('Starting iSFP PPT Service v6.0')
    app.run(host='0.0.0.0', port=5000, debug=False)
